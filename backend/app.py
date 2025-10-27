from fastapi import FastAPI, HTTPException, UploadFile, status, File, Depends, Form
from pymongo import MongoClient
from bson import ObjectId
import boto3
from typing import List
from Models.model import CSVNodeQueryRequest, CSVNodeQueryResponse, Flow
from Models.model import PDFNodeQueryRequest
from Models.model import PDFNodeQueryResponse
from Models.model import TXTNodeQueryRequest
from Models.model import TXTNodeQueryResponse
from Models.model import MDNodeQueryRequest
from Models.model import MDNodeQueryResponse
from Models.model import HTMLNodeQueryRequest
from Models.model import HTMLNodeQueryResponse
from Models.model import DOCXNodeQueryRequest
from Models.model import DOCXNodeQueryResponse
from Models.model import PPTXNodeQueryRequest
from Models.model import PPTXNodeQueryResponse
from Models.model import ImgNodeQueryRequest
from Models.model import ImgNodeQueryResponse
from Models.model import AudioNodeQueryRequest
from Models.model import AudioNodeQueryResponse
from Models.model import YoutubeNodeQueryRequest
from Models.model import YoutubeNodeQueryResponse
from Models.model import VideoNodeQueryRequest
from Models.model import VideoNodeQueryResponse
from Models.model import WebNodeQueryRequest
from Models.model import WebNodeQueryResponse
from Models.model import SQLComponentRequest
from Models.model import SQLComponentResponse
from Models.model import SQLNodeQueryRequest
from Models.model import SQLNodeQueryResponse
from Models.model import ComponentFollowUpQueryRequest
from Models.model import ComponentFollowUpQueryResponse
from Models.model import MultipleQuestionAnswerQueryRequest
from Models.model import MultipleQuestionAnswerQueryResponse
from Models.model import FlowSummarizeRequest
from Models.model import FlowSummarizeResponse
from pymongo.mongo_client import MongoClient
from botocore.exceptions import ClientError
from hashlib import sha256
from io import BytesIO
from trp import Document
import time
import pandas as pd
from langchain_experimental.text_splitter import SemanticChunker
from langchain_openai.embeddings import OpenAIEmbeddings
from langchain_chroma import Chroma
from langchain_community.document_loaders import UnstructuredMarkdownLoader
import chromadb
from uuid import uuid4
from langchain_core.documents import Document as LangDocument
from langchain_openai import ChatOpenAI
from langchain_core.prompts import PromptTemplate
import datetime
from fastapi.middleware.cors import CORSMiddleware
import pandas as pd
import sqlite3
from vanna.openai import OpenAI_Chat
from vanna.chromadb import ChromaDB_VectorStore
from typing import List
from crawl4ai import *
import os
import base64
import openai
import pypdfium2 as pdfium
from langchain.prompts import PromptTemplate
from langchain.chains import LLMChain
from langchain.chains.combine_documents.stuff import StuffDocumentsChain
from langchain.chains.combine_documents.reduce import ReduceDocumentsChain
from langchain.chains import MapReduceDocumentsChain
import json
# Conditional imports for PDF processing
try:
    from unstructured.partition.pdf import partition_pdf
    from unstructured.documents.elements import (
        Text,
        Title,
        NarrativeText,
        ListItem,
        Header,
    )
    UNSTRUCTURED_AVAILABLE = True
except ImportError as e:
    print(f"Warning: unstructured package not available: {e}")
    print("PDF processing with unstructured will be disabled.")
    UNSTRUCTURED_AVAILABLE = False

try:
    import camelot
    CAMELOT_AVAILABLE = True
except ImportError as e:
    print(f"Warning: camelot package not available: {e}")
    print("PDF table extraction with camelot will be disabled.")
    CAMELOT_AVAILABLE = False

import re
import google.generativeai as genai
import vertexai
from vertexai.generative_models import GenerativeModel, Part
from google.oauth2 import service_account
import json as json_lib
import io
import markdown
from bs4 import BeautifulSoup
from docx import Document
from pptx import Presentation
import tiktoken
import traceback
from dotenv import load_dotenv

load_dotenv()

mongo_db_url = os.getenv("mongo_db_url")
openai_api_key_str = os.getenv("openai_api_key")
gemini_api_key_str = os.getenv("gemini_api_key")
gcp_project_id_str = os.getenv("gcp_project_id")
aws_access_key_id_str = os.getenv("aws_access_key_id")
aws_secret_access_key_str = os.getenv("aws_secret_access_key")
bucket_name = os.getenv("bucket_name")

# New environment variables for Bedrock configuration
use_bedrock_only = os.getenv("USE_BEDROCK_ONLY", "false").lower() == "true"
aws_region = os.getenv("AWS_REGION", "us-east-1")

# Initialize Bedrock client
bedrock_runtime_client = boto3.client(
    'bedrock-runtime',
    region_name=aws_region
)

# Initialize Google models only if not using Bedrock exclusively
if not use_bedrock_only:
    try:
        credentials = service_account.Credentials.from_service_account_file(
            "./ai-interview-poc-2b5cf8540f16.json"
        )
        
        vertexai.init(
            project=gcp_project_id_str, credentials=credentials, location="us-central1"
        )
        
        model_vertexai = GenerativeModel("gemini-2.0-flash")
        
        genai.configure(api_key=gemini_api_key_str)
        
        model = genai.GenerativeModel("gemini-2.0-flash")
    except Exception as e:
        print(f"Warning: Could not initialize Google models: {e}")
        model_vertexai = None
        model = None
else:
    model_vertexai = None
    model = None

connection = sqlite3.connect("sqlite_data.db")

GPT_4O_MAX_TOKENS = 128000

UPLOAD_DIR = "uploaded_pdfs"

if not use_bedrock_only:
    openai.api_key = openai_api_key_str

# Bedrock helper functions
def get_aws_account_id():
    """Get AWS account ID"""
    try:
        sts_client = boto3.client('sts')
        return sts_client.get_caller_identity()['Account']
    except Exception:
        return None

def analyze_content_with_pegasus(s3_uri, prompt, temperature=0):
    """Analyze content using AWS Bedrock Pegasus model"""
    try:
        account_id = get_aws_account_id()
        
        if not account_id:
            return "Error: Unable to retrieve AWS account ID"
        
        request_body = {
            "inputPrompt": prompt,
            "mediaSource": {
                "s3Location": {
                    "uri": s3_uri,
                    "bucketOwner": account_id
                }
            },
            "temperature": temperature
        }
        
        response = bedrock_runtime_client.invoke_model(
            modelId="us.twelvelabs.pegasus-1-2-v1:0",
            body=json_lib.dumps(request_body),
            contentType="application/json",
            accept="application/json"
        )
        
        response_body = json_lib.loads(response["body"].read())
        return response_body.get("message", "No response received")
    except Exception as e:
        return f"Error analyzing content with Pegasus: {str(e)}"

def generate_text_with_claude(prompt, max_tokens=4000):
    """Generate text using AWS Bedrock Claude Sonnet 3.5"""
    try:
        request_body = {
            "anthropic_version": "bedrock-2023-05-31",
            "max_tokens": max_tokens,
            "messages": [
                {
                    "role": "user",
                    "content": prompt
                }
            ]
        }
        
        response = bedrock_runtime_client.invoke_model(
            modelId="anthropic.claude-3-5-sonnet-20241022-v2:0",
            body=json_lib.dumps(request_body),
            contentType="application/json",
            accept="application/json"
        )
        
        response_body = json_lib.loads(response["body"].read())
        return response_body["content"][0]["text"]
    except Exception as e:
        return f"Error generating text with Claude: {str(e)}"

def analyze_image_with_claude(image_base64, prompt, mime_type="image/jpeg"):
    """Analyze image using AWS Bedrock Claude Sonnet 3.5"""
    try:
        request_body = {
            "anthropic_version": "bedrock-2023-05-31",
            "max_tokens": 4000,
            "messages": [
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "image",
                            "source": {
                                "type": "base64",
                                "media_type": mime_type,
                                "data": image_base64
                            }
                        },
                        {
                            "type": "text",
                            "text": prompt
                        }
                    ]
                }
            ]
        }
        
        response = bedrock_runtime_client.invoke_model(
            modelId="anthropic.claude-3-5-sonnet-20241022-v2:0",
            body=json_lib.dumps(request_body),
            contentType="application/json",
            accept="application/json"
        )
        
        response_body = json_lib.loads(response["body"].read())
        return response_body["content"][0]["text"]
    except Exception as e:
        return f"Error analyzing image with Claude: {str(e)}"

class SQLBot(ChromaDB_VectorStore, OpenAI_Chat):
    def __init__(self, config=None):
        ChromaDB_VectorStore.__init__(self, config=config)
        OpenAI_Chat.__init__(self, config=config)


class CSVBot(ChromaDB_VectorStore, OpenAI_Chat):
    def __init__(self, config=None):
        ChromaDB_VectorStore.__init__(self, config=config)
        OpenAI_Chat.__init__(self, config=config)


sqlBot = SQLBot(
    config={
        "api_key": openai_api_key_str,
        "model": "gpt-4o",
        "temperature": 0,
        "path": "./SQLVectorStore",
        "client": "persistent",
        "n_results": 1,
    }
)

csvBot = CSVBot(
    config={
        "api_key": openai_api_key_str,
        "model": "gpt-4o",
        "temperature": 0,
        "path": "./CSVVectorStore",
        "client": "persistent",
        "n_results": 1,
    }
)

sqlBot.connect_to_sqlite("sqlite_data.db")

csvBot.connect_to_sqlite("csv_data.db")


df_ddl = sqlBot.run_sql("SELECT type, sql FROM sqlite_master WHERE sql IS NOT NULL")
# print(df_ddl)

for ddl in df_ddl['sql'].to_list():
    sqlBot.train(ddl=ddl)

# training_data = sqlBot.get_training_data()
# print(training_data)

app = FastAPI()

origins = ["*"]

app.add_middleware(
    CORSMiddleware,
    allow_origins=origins,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


client = MongoClient(mongo_db_url)

try:
    client.admin.command("ping")
    print("Pinged your deployment. You successfully connected to MongoDB!")
except Exception as e:
    print(e)

db = client["MindMap"]
flow_collection = db["flows"]
component_collection = db["components"]
node_collection = db["nodes"]

# AWS S3 setup
s3_client = boto3.client(
    "s3",
    aws_access_key_id=aws_access_key_id_str,
    aws_secret_access_key=aws_secret_access_key_str,
    region_name="ap-south-1",
)


embedding_function = OpenAIEmbeddings(
    model="text-embedding-ada-002", api_key=openai_api_key_str
)

persistent_client = chromadb.PersistentClient()

vector_store_from_client = Chroma(
    client=persistent_client,
    collection_name="pdfs",
    embedding_function=embedding_function,
)

PDFCollection = persistent_client.get_collection("pdfs")

text_splitter = SemanticChunker(
    OpenAIEmbeddings(model="text-embedding-ada-002", api_key=openai_api_key_str)
)

# Initialize LLM based on configuration
if use_bedrock_only:
    # Create a custom LLM wrapper for Claude
    class ClaudeLLM:
        def invoke(self, input_data):
            if isinstance(input_data, dict) and "text" in input_data:
                prompt = input_data["text"]
            elif hasattr(input_data, 'text'):
                prompt = input_data.text
            else:
                prompt = str(input_data)
            
            response_text = generate_text_with_claude(prompt)
            
            # Return object with content attribute to match OpenAI interface
            class Response:
                def __init__(self, content):
                    self.content = content
            
            return Response(response_text)
    
    llm = ClaudeLLM()
else:
    llm = ChatOpenAI(model="gpt-4o", api_key=openai_api_key_str)

def calculate_file_hash(file):
    hasher = sha256()

    if isinstance(file, bytes):
        file = BytesIO(file)

    while chunk := file.read(8192):
        hasher.update(chunk)

    file.seek(0)
    return hasher.hexdigest()


def validate_dataframe(df):
    try:
        if not isinstance(df, list) and not all(isinstance(item, dict) for item in df):
            return []
        else:
            return df

    except ValueError:
        return []


def upload_to_s3(file_bytes, bucket, key):
    try:
        s3_client.put_object(Bucket=bucket, Key=key, Body=file_bytes)
    except ClientError as e:
        raise HTTPException(status_code=500, detail=f"S3 Upload Error: {str(e)}")


def extract_text_and_tables(key):
    client = boto3.client(
        "textract",
        aws_access_key_id=aws_access_key_id_str,
        aws_secret_access_key=aws_secret_access_key_str,
        region_name="ap-south-1",
    )

    response = client.start_document_analysis(
        DocumentLocation={"S3Object": {"Bucket": bucket_name, "Name": key}},
        FeatureTypes=["TABLES", "FORMS"],
    )

    print(response)

    job_id = response["JobId"]
    print("job_id")
    while True:
        response = client.get_document_analysis(JobId=job_id)
        status = response["JobStatus"]
        if status in ["SUCCEEDED", "FAILED"]:
            break
        time.sleep(5)

    if status == "FAILED":
        raise Exception("Document analysis failed")

    all_blocks = response["Blocks"]
    next_token = response.get("NextToken", None)
    print(all_blocks)

    while next_token:
        response = client.get_document_analysis(JobId=job_id, NextToken=next_token)
        all_blocks.extend(response["Blocks"])
        next_token = response.get("NextToken", None)

    response["Blocks"] = all_blocks

    doc = Document(response)
    lines = [line.text for page in doc.pages for line in page.lines if line.text]
    tables = []
    key_values = {}

    for page in doc.pages:
        for table in page.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text if cell.text else "" for cell in row.cells]
                table_data.append(row_data)
            df = pd.DataFrame(table_data)
            tables.append(df)

        for field in page.form.fields:
            key = field.key.text if field.key and field.key.text else ""
            value = field.value.text if field.value and field.value.text else ""
            key_values[key] = value

    date = next((line for line in lines if "Date" in line), None)
    return lines, tables, key_values, date


def sanitize_path(path):
    """Sanitize the path to remove any invalid characters."""
    return re.sub(r"[^\w\s-]", "_", path).strip()


def camelot_pdf_processing(flow_id, file, flow_type):
    try:
        # Check if required packages are available
        if not UNSTRUCTURED_AVAILABLE:
            raise HTTPException(
                status_code=500, 
                detail="PDF processing unavailable: unstructured package not installed. Please install system dependencies."
            )
        
        if not CAMELOT_AVAILABLE:
            print("Warning: Camelot not available, table extraction will be skipped")
        
        try:
        sanitized_flow_id = sanitize_path(flow_id)
        sanitized_filename = file.filename

        flow_dir = os.path.join(UPLOAD_DIR, sanitized_flow_id)

        # Ensure the directory exists
        os.makedirs(flow_dir, exist_ok=True)

        # Create the full file path
        file_path = os.path.join(flow_dir, sanitized_filename)

        print(file_path)

        # Save the uploaded file
        with open(file_path, "wb") as f:
            # Read the file content and write it to disk
            f.write(file.file.read())

        file.file.seek(0)
        file_bytes = file.file.read()

        print(f"File saved at: {file_path}")

        print(file_bytes)

        file_hash = calculate_file_hash(file_bytes)
        print(file_hash)

        existing_component = component_collection.find_one(
            {"file_hash": file_hash, "flow_id": ObjectId(flow_id)}
        )
        if existing_component:
            print("File already exists in the system")
            raise HTTPException(
                status_code=400, detail="File already exists in the system."
            )

        # Now, upload the file to S3
        folder = f"uploads/{flow_id}/"
        s3_key = folder + sanitized_filename

        # Pass the local file path to upload_to_s3
        upload_to_s3(file_bytes, bucket_name, s3_key)
        print("File uploaded to S3")

        ocr_config = {
            "languages": ["eng"],
            "strategy": "fast",
        }
        elements = partition_pdf(filename=file_path, **ocr_config)

        content_data = []

        # Extract content (text, title, etc.) from the PDF
        for element in elements:
            if isinstance(element, (Text, Title, NarrativeText, ListItem, Header)):
                content_data.append(
                    {"data": element.text, "page_number": element.metadata.page_number}
                )

        result_tbl_list = []
        
        # Extract tables from the PDF using camelot if available
        if CAMELOT_AVAILABLE:
            try:
                tables = camelot.read_pdf(
                    file_path,
                    pages="all",
                    flavor="stream",
                    edge_tol=900,
                    row_tol=6,
                    strip_text="\n",
                )

                print(f"Total tables found: {len(tables)}")

                # Extract tables from the PDF
                for i, table in enumerate(tables):
                    page_number = (
                        table.page
                    )  # Get the page number from which the table was extracted
                    print(f"Table {i + 1} extracted from page: {page_number}\n")
                    print(table.df)
                    result_tbl_list.append(
                        {"data": table.df.to_string(index=False), "page_number": table.page}
                    )
            except Exception as e:
                print(f"Warning: Camelot table extraction failed: {e}")
                result_tbl_list = []
        else:
            print("Camelot not available, skipping table extraction")

        # Combine content_data and result_tbl_list based on the page number
        combined_data = []

        for content_item in content_data:
            page_number = content_item["page_number"]
            # Check if a matching table exists for the current page number
            matching_table = next(
                (
                    table
                    for table in result_tbl_list
                    if table["page_number"] == page_number
                ),
                None,
            )

            # Store content along with the page number and table (if exists)
            if matching_table:
                combined_content = (
                    f"{content_item['data']}\n\nTable:\n{matching_table['data']}"
                )
                combined_data.append(
                    {"page_number": page_number, "content": combined_content}
                )
            else:
                combined_data.append(
                    {"page_number": page_number, "content": content_item["data"]}
                )

        print("-==============-")

        # Print the final combined data
        print(combined_data)

        combined_data_str = json.dumps(combined_data, indent=4)

        chunks = do_semantic_chunking(combined_data_str)
        print(chunks)

        summary = process_pdf_summary(chunks)
        
        if flow_type == "manual":

            component_metadata = {
                "flow_id": ObjectId(flow_id),
                "name": file.filename,
                "file_hash": file_hash,
                "size": len(file_bytes),
                "s3_path": s3_key,
                "type": "pdf",
                "processing_type": "custom",
                "summary": summary,
            }

            component_id = component_collection.insert_one(component_metadata).inserted_id
            EmbeddingsDocuments = []
            for i in range(len(chunks)):
                metadata = {
                    "component_id": str(component_id),
                    "file_name": file.filename,
                    "flow_id": flow_id,
                }
                EmbeddingsDocuments.append(
                    LangDocument(page_content=chunks[i].page_content, metadata=metadata)
                )
            uuids = [str(uuid4()) for _ in range(len(EmbeddingsDocuments))]

            vector_store_from_client.add_documents(documents=EmbeddingsDocuments, ids=uuids)
            return {"component_id": str(component_id), "type": "pdf"}
        
        else:
                    
            template = """You are tasked with generating a JSON mind map for given summary of the pdf document and that should be compatible with React Flow for rendering a flow diagram. The mind map should adhere to the following rules:

                1. **Node Types:**
                - There will always be one `dataSource` node, which serves as the root of the flow.
                - There will be a maximum of 5 `response` nodes.

                2. **Node Relationships:**
                - The `dataSource` node should be connected to all `response` nodes.
                - `response` nodes may also connect to each other if it improves the logical flow or visualization.

                3. **Node Properties:**
                - Each node should have:
                    - `id` (unique identifier of 12 or 24 digit unique uuid or nanoid)
                    - `type` (`dataSource` or `response`)
                    - `position` (coordinates in the form {{ "x": <number>, "y": <number> }} for layout)
                    - `measured` (an object defining width and height):
                        {{
                            "width": <number>,
                            "height": <number>
                        }}
                    - `targetPosition` (position of the target connection, default to `"left"`)
                    - `sourcePosition` (position of the source connection, default to `"right"`)
                    - `selected` (boolean, default to `false`)
                    - `deletable` (boolean, default to `true` for `response` and `false` for `dataSource`)

                4. **Node Data Format:**
                - `dataSource` Node:
                    - `data` contains the following properties:
                        {{
                            "prompt": "<data source description>",
                            "name": "pdf", !!!DOESN"T CHANGES 
                            "content": "<file name or content>",
                            "flow_id": "{flow_id}",
                            "file": "{filename}"  // Empty object or file metadata
                        }}

                - `response` Node:
                    - `data` contains nested properties:
                        {{
                            "id": "<unique identifier of 12 or 24 digit unique uuid or nanoid>",
                            "type": "MDNode | WebNode | MultipleQA | other",
                            "data": {{
                                "question": "<question text, if applicable>",
                                "summ": "<summary or answer>",
                                "df": [],
                                "graph": "",
                                "flow_id": "{flow_id}",
                                "component_id": "<component reference ID - unique identifier of 12 or 24 digit unique uuid or nanoid>",
                                "component_type": "pdf"
                            }}
                        }}

                5. **Connections:**
                - Connections between nodes should be represented by edges, with the following format:
                    - `id` (unique identifier for the edge)
                    - `source` (ID of the source node)
                    - `target` (ID of the target node)
                    - `type` (optional, defaults to `default`)
                    - 'animated' !!WILL ALWAYS BE TRUE

                6. **Viewport Configuration:**
                - Include a `viewport` object that specifies:
                    - `x` (horizontal position of the viewport)
                    - `y` (vertical position of the viewport)
                    - `zoom` (zoom level for initial rendering)
                    
                Here is the PDF summary for which you need to generate the mind map:
                {summary_pdf}

                ### Additional Considerations:
                - Ensure that the node positions are distributed properly to avoid overlap.
                - If fewer than 5 `response` nodes are required, adjust accordingly.
                - Prioritize connecting `response` nodes where it adds logical structure to the flow.

                ### IMPORTANT:
                - **RETURN ONLY THE VALID JSON OBJECT AND NO ADDITIONAL COMMENTS**.
                - Do **not** include any explanations, text, or additional information.
                - Maintain the format with double curly braces `{{` and `}}` as shown in the format.
                """

            prompt = PromptTemplate.from_template(template)

            lm_chain = prompt | llm
            
            answer = lm_chain.invoke(
                    {"summary_pdf": summary, "flow_id": flow_id, "filename": file.filename}
            )

            responseList = answer.content

            print(responseList)

            response_json =  response_json.replace("```json", "").replace("```", "").replace("\n", "").strip()
            
            print(response_json)
            
            component_metadata = {
                "flow_id": ObjectId(flow_id),
                "name": file.filename,
                "type": "pdf",
                "processing_type": "gpt",
                "mindmap_json": response_json,
            }

            component_id = component_collection.insert_one(component_metadata).inserted_id

            return {
                "component_id": str(component_id),
                "type": "pdf",
                "mindmap_json": response_json,
                "flow_type": "automatic"
            }
            
    except Exception as e:
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=e)


def extract_text_from_upload(file: UploadFile) -> str:
    file.file.seek(0)
    content = file.file.read()
    extension = file.filename.split(".")[-1].lower()

    if extension == "txt":
        return content.decode("utf-8", errors="ignore")

    elif extension == "md":
        html = markdown.markdown(content.decode("utf-8", errors="ignore"))
        return BeautifulSoup(html, "html.parser").get_text()

    elif extension == "html":
        return BeautifulSoup(content, "html.parser").get_text()

    elif extension == "docx":
        doc = Document(io.BytesIO(content))
        return "\n".join([p.text for p in doc.paragraphs])

    elif extension == "pptx":
        prs = Presentation(io.BytesIO(content))
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    else:
        raise ValueError(f"Unsupported file type: {extension}")


def is_within_gpt4o_token_limit(file: UploadFile) -> bool:
    try:
        text = extract_text_from_upload(file)
        encoding = tiktoken.get_encoding("cl100k_base")
        token_count = len(encoding.encode(text))
        print("token_count : ", token_count)
        return token_count <= GPT_4O_MAX_TOKENS
    except Exception as e:
        print(f"Error checking token count: {e}")
        return False


def process_pdf_summary(chunks):
    map_template = """Write a concise summary of the following content:
                    {content}
                    Summary:
                    """
    map_prompt = PromptTemplate.from_template(map_template)
    map_chain = LLMChain(prompt=map_prompt, llm=llm)

    reduce_template = """The following is a set of summaries:
                        {doc_summaries}
                        Summarize the above summaries with all the key details.
                        Summary:"""
    reduce_prompt = PromptTemplate.from_template(reduce_template)
    reduce_chain = LLMChain(prompt=reduce_prompt, llm=llm)

    stuff_chain = StuffDocumentsChain(
        llm_chain=reduce_chain, document_variable_name="doc_summaries"
    )

    reduce_documents_chain = ReduceDocumentsChain(
        combine_documents_chain=stuff_chain,
    )

    map_reduce_chain = MapReduceDocumentsChain(
        llm_chain=map_chain,
        document_variable_name="content",
        reduce_documents_chain=reduce_documents_chain,
    )

    small_chunks = chunks[:24]

    summary = map_reduce_chain.run(small_chunks)
    print(summary)
    return summary


def use_aws_textract(file, flow_id, flow_type):
    if not file.filename.endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Only PDF files are allowed.")
    file_bytes = file.file.read()
    file_hash = calculate_file_hash(file_bytes)
    print(file_hash)
    existing_component = component_collection.find_one(
        {"file_hash": file_hash, "flow_id": ObjectId(flow_id)}
    )
    if existing_component:
        raise HTTPException(
            status_code=400, detail="File already exists in the system."
        )
    file_name = file.filename
    folder = f"uploads/{flow_id}/"
    s3_key = folder + file_name
    upload_to_s3(file_bytes, bucket_name, s3_key)
    print("uploaded")
    data_for_chunking = "Data : "
    try:
        lines, tables, key_values, date = extract_text_and_tables(s3_key)
        for i, table in enumerate(tables):
            print(f"Table {i+1}:\n")
            print(table.to_string() + "\n\n")
            data_for_chunking = data_for_chunking + f"Table {i+1}:\n"
            data_for_chunking = data_for_chunking + table.to_string() + "\n\n"
        print("Extracted Key-Value Pairs:\n")
        for key, value in key_values.items():
            print(f"{key}: {value}\n")
            data_for_chunking = data_for_chunking + f"{key}: {value}\n"
        if date:
            print(f"Date: {date}\n")
            data_for_chunking = data_for_chunking + f"Date: {date}\n"
    except Exception as e:
        print(e)
        raise HTTPException(status_code=500, detail=f"Text extraction error: {str(e)}")
    chunks = do_semantic_chunking(data_for_chunking)
    print(chunks)

    summary = process_pdf_summary(chunks)
    
    if flow_type == "manual":
        
        component_metadata = {
            "flow_id": ObjectId(flow_id),
            "name": file.filename,
            "file_hash": file_hash,
            "size": len(file_bytes),
            "s3_path": s3_key,
            "type": "pdf",
            "processing_type": "aws",
            "summary": summary,
        }

        component_id = component_collection.insert_one(component_metadata).inserted_id
        EmbeddingsDocuments = []
        for i in range(len(chunks)):
            metadata = {
                "component_id": str(component_id),
                "file_name": file.filename,
                "flow_id": flow_id,
            }
            EmbeddingsDocuments.append(
                LangDocument(page_content=chunks[i].page_content, metadata=metadata)
            )
        uuids = [str(uuid4()) for _ in range(len(EmbeddingsDocuments))]

        vector_store_from_client.add_documents(documents=EmbeddingsDocuments, ids=uuids)
        return {"component_id": str(component_id), "type": "pdf"}
    
    else:
        
        template = """You are tasked with generating a JSON mind map for given summary of the pdf document and that should be compatible with React Flow for rendering a flow diagram. The mind map should adhere to the following rules:

                1. **Node Types:**
                - There will always be one `dataSource` node, which serves as the root of the flow.
                - There will be a maximum of 5 `response` nodes.

                2. **Node Relationships:**
                - The `dataSource` node should be connected to all `response` nodes.
                - `response` nodes may also connect to each other if it improves the logical flow or visualization.

                3. **Node Properties:**
                - Each node should have:
                    - `id` (unique identifier of 12 or 24 digit unique uuid or nanoid)
                    - `type` (`dataSource` or `response`)
                    - `position` (coordinates in the form {{ "x": <number>, "y": <number> }} for layout)
                    - `measured` (an object defining width and height):
                        {{
                            "width": <number>,
                            "height": <number>
                        }}
                    - `targetPosition` (position of the target connection, default to `"left"`)
                    - `sourcePosition` (position of the source connection, default to `"right"`)
                    - `selected` (boolean, default to `false`)
                    - `deletable` (boolean, default to `true` for `response` and `false` for `dataSource`)

                4. **Node Data Format:**
                - `dataSource` Node:
                    - `data` contains the following properties:
                        {{
                            "prompt": "<data source description>",
                            "name": "pdf", !!!DOESN"T CHANGES 
                            "content": "<file name or content>",
                            "flow_id": "{flow_id}",
                            "file": "{filename}"  // Empty object or file metadata
                        }}

                - `response` Node:
                    - `data` contains nested properties:
                        {{
                            "id": "<unique identifier of 12 or 24 digit unique uuid or nanoid>",
                            "type": "MDNode | WebNode | MultipleQA | other",
                            "data": {{
                                "question": "<question text, if applicable>",
                                "summ": "<summary or answer>",
                                "df": [],
                                "graph": "",
                                "flow_id": "{flow_id}",
                                "component_id": "<component reference ID - unique identifier of 12 or 24 digit unique uuid or nanoid>",
                                "component_type": "pdf"
                            }}
                        }}

                5. **Connections:**
                - Connections between nodes should be represented by edges, with the following format:
                    - `id` (unique identifier for the edge)
                    - `source` (ID of the source node)
                    - `target` (ID of the target node)
                    - `type` (optional, defaults to `default`)
                    - 'animated' !!WILL ALWAYS BE TRUE

                6. **Viewport Configuration:**
                - Include a `viewport` object that specifies:
                    - `x` (horizontal position of the viewport)
                    - `y` (vertical position of the viewport)
                    - `zoom` (zoom level for initial rendering)
                    
                Here is the PDF summary for which you need to generate the mind map:
                {summary_pdf}

                ### Additional Considerations:
                - Ensure that the node positions are distributed properly to avoid overlap.
                - If fewer than 5 `response` nodes are required, adjust accordingly.
                - Prioritize connecting `response` nodes where it adds logical structure to the flow.

                ### IMPORTANT:
                - **RETURN ONLY THE VALID JSON OBJECT AND NO ADDITIONAL COMMENTS**.
                - Do **not** include any explanations, text, or additional information.
                - Maintain the format with double curly braces `{{` and `}}` as shown in the format.
                """

        prompt = PromptTemplate.from_template(template)

        lm_chain = prompt | llm
        
        answer = lm_chain.invoke(
                {"summary_pdf": summary, "flow_id": flow_id, "filename": file.filename}
        )

        responseList = answer.content

        print(responseList)

        response_json =  response_json.replace("```json", "").replace("```", "").replace("\n", "").strip()
        
        print(response_json)
        
        component_metadata = {
            "flow_id": ObjectId(flow_id),
            "name": file.filename,
            "type": "pdf",
            "processing_type": "gpt",
            "mindmap_json": response_json,
        }

        component_id = component_collection.insert_one(component_metadata).inserted_id

        return {
            "component_id": str(component_id),
            "type": "pdf",
            "mindmap_json": response_json,
            "flow_type": "automatic"
        }


def do_semantic_chunking(docs):
    documents = text_splitter.create_documents([docs])
    print("Number of chunks created: ", len(documents))
    for i in range(len(documents)):
        print()
        print(f"CHUNK : {i+1}")
        print(documents[i].page_content)
    return documents


def get_relevant_passage(query: str, flow_id: str, component_id: str, n_results: int):
    results = vector_store_from_client.similarity_search(
        query=query,
        k=2,
        filter={"$and": [{"component_id": component_id}, {"flow_id": flow_id}]},
    )
    return [doc.page_content for doc in results]


def fetch_question_answer_from_node_collection(parent_id: str, flow_id: str):
    try:
        record = node_collection.find_one(
            {
                "_id": ObjectId(parent_id),
                "flow_id": ObjectId(flow_id),
                "is_delete": "false",
            }
        )

        print("Fetched record:", record)

        if not record:
            return None, None

        print("Question value:", record.get("question"))

        question = record.get("question", None)
        answer = None

        if record["type"] == "pdf":
            answer = "Answer: " + str(record.get("summ", "Answer not found."))
            df_list = record.get("df", [])
            df_string = " | ".join([str(item) for item in df_list])

            answer += " DataFrame: " + df_string

        elif record["type"] == "txt":
            answer = "Answer: " + str(record.get("summ", "Answer not found."))
            df_list = record.get("df", [])
            df_string = " | ".join([str(item) for item in df_list])

            answer += " DataFrame: " + df_string

        elif record["type"] == "md":
            answer = "Answer: " + str(record.get("summ", "Answer not found."))
            df_list = record.get("df", [])
            df_string = " | ".join([str(item) for item in df_list])

            answer += " DataFrame: " + df_string

        elif record["type"] == "html":
            answer = "Answer: " + str(record.get("summ", "Answer not found."))
            df_list = record.get("df", [])
            df_string = " | ".join([str(item) for item in df_list])

            answer += " DataFrame: " + df_string

        elif record["type"] == "docx":
            answer = "Answer: " + str(record.get("summ", "Answer not found."))
            df_list = record.get("df", [])
            df_string = " | ".join([str(item) for item in df_list])

            answer += " DataFrame: " + df_string

        elif record["type"] == "pptx":
            answer = "Answer: " + str(record.get("summ", "Answer not found."))
            df_list = record.get("df", [])
            df_string = " | ".join([str(item) for item in df_list])

            answer += " DataFrame: " + df_string

        elif record["type"] == "image":
            answer = "Answer: " + str(record.get("summ", "Answer not found."))
            df_list = record.get("df", [])

            df_string = " | ".join([str(item) for item in df_list])

            answer += " DataFrame: " + df_string

        elif record["type"] == "audio":
            answer = "Answer: " + str(record.get("summ", "Answer not found."))
            df_list = record.get("df", [])
            df_string = " | ".join([str(item) for item in df_list])

            answer += " DataFrame: " + df_string

        elif record["type"] == "youtube":
            answer = "Answer: " + str(record.get("summ", "Answer not found."))
            df_list = record.get("df", [])
            df_string = " | ".join([str(item) for item in df_list])

            answer += " DataFrame: " + df_string

        elif record["type"] == "video":
            answer = "Answer: " + str(record.get("summ", "Answer not found."))
            df_list = record.get("df", [])
            df_string = " | ".join([str(item) for item in df_list])

            answer += " DataFrame: " + df_string

        elif record["type"] == "sql":
            answer = "Answer: " + str(record.get("summ", "Answer not found."))
            df_list = record.get("df", [])
            df_string = " | ".join([str(item) for item in df_list])

            answer += " DataFrame: " + df_string

        elif record["type"] == "csv":
            answer = "Answer: " + str(record.get("summ", "Answer not found."))
            df_list = record.get("df", [])
            df_string = " | ".join([str(item) for item in df_list])
            
            answer += " DataFrame: " + df_string

        elif record["type"] == "web":
            answer = "Answer: " + str(record.get("summ", "Answer not found."))
            df_list = record.get("df", [])
            df_string = " | ".join([str(item) for item in df_list])
            
            answer += " DataFrame: " + df_string

        elif record["type"] == "MultipleQA":
            answer = "Answer: " + str(record.get("summ", "Answer not found."))
            df_list = record.get("df", [])
            df_string = " | ".join([str(item) for item in df_list])
            
            answer += " DataFrame: " + df_string

        print("Answer:", answer)

        return question, answer
    except Exception as e:
        traceback.print_exc()
        print(f"Error: {e}")
        return None, None


@app.post("/create-flow/")
def create_flow(flow: dict):
    try:
        flow_data = {"flow_name": flow.get("flow_name"), "flow_json": "", "summary": "", "flow_type": flow.get("flow_type")}
        flow_id = flow_collection.insert_one(flow_data).inserted_id
        flow_type = flow.get("flow_type")
        return {"flow_id": str(flow_id), "flow_type": str(flow_type)}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error creating flow: {str(e)}")


@app.delete("/delete-flow/{flow_id}")
def delete_flow(flow_id: str):
    try:
        flow_object_id = ObjectId(flow_id)

        components = component_collection.find({"flow_id": flow_object_id})
        for component in components:
            component_id = component["_id"]
            node_collection.delete_many({"component_id": component_id})
        component_collection.delete_many({"flow_id": flow_object_id})

        flow_collection.delete_one({"_id": flow_object_id})

        return {"status": "success"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error deleting flow: {str(e)}")


@app.get("/flows/", response_model=List[Flow])
def list_flows():
    flows = flow_collection.find()
    return [
        {
            "flow_id": str(flow["_id"]),
            "flow_name": flow["flow_name"],
            "flow_json": flow["flow_json"],
            "summary": flow["summary"],
            "flow_type": flow["flow_type"]
        }
        for flow in flows
    ]


@app.put("/flow-update/")
def update_flow(update_data: Flow):
    try:
        print(update_data)
        result = flow_collection.update_one(
            {"_id": ObjectId(update_data.flow_id)},
            {
                "$set": {
                    "flow_name": update_data.flow_name,
                    "flow_json": update_data.flow_json,
                    "flow_type": update_data.flow_type,
                    "summary": update_data.summary,
                }
            },
        )
        print(result)

        if result.matched_count == 0:
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND, detail="Flow not found"
            )

        return {
            "flow_id": str(update_data.flow_id),
            "message": "Flow updated successfully",
        }

    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"An error occurred: {str(e)}",
        )


def get_summary_from_openai(file: UploadFile, flow_id: str, flow_type: str):
    file.file.seek(0)
    file_bytes = file.file.read()
    file_extension = file.filename.split(".")[-1].lower()

    if len(file_bytes) == 0:
        raise ValueError("The uploaded file is actually empty!")

    assistant = openai.beta.assistants.create(
        name="Summarize agent",
        instructions="Your task is to only summarize the document",
        model="gpt-4o",
        tools=[{"type": "file_search"}],
    )
    vector_store = openai.beta.vector_stores.create(name=f"{file_extension}_{flow_id}")

    assistant = openai.beta.assistants.update(
        assistant_id=assistant.id,
        tool_resources={"file_search": {"vector_store_ids": [vector_store.id]}},
    )

    if file_extension == "pdf":
        mime_type = "application/pdf"
        messages_file = openai.files.create(
            file=(file.filename, file_bytes, mime_type), purpose="assistants"
        )
    elif file_extension == "csv":
        mime_type = "application/json"
        df = pd.read_csv(BytesIO(file_bytes))
        json_data = df.to_dict(orient="records")
        json_str = json.dumps(json_data)
        messages_file = openai.files.create(
            file=("data.json", json_str, mime_type), purpose="assistants"
        )
    elif file_extension == "txt":
        mime_type = "text/plain"
        messages_file = openai.files.create(
            file=(file.filename, file_bytes, mime_type), purpose="assistants"
        )
    elif file_extension == "docx":
        mime_type = (
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        )
        messages_file = openai.files.create(
            file=(file.filename, file_bytes, mime_type), purpose="assistants"
        )
    elif file_extension == "html":
        mime_type = "text/html"
        messages_file = openai.files.create(
            file=(file.filename, file_bytes, mime_type), purpose="assistants"
        )
    elif file_extension == "md":
        mime_type = "text/markdown"
        messages_file = openai.files.create(
            file=(file.filename, file_bytes, mime_type), purpose="assistants"
        )
    elif file_extension == "pptx":
        mime_type = (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
        messages_file = openai.files.create(
            file=(file.filename, file_bytes, mime_type), purpose="assistants"
        )
    else:
        raise ValueError("Unsupported file type")

    # Create thread for summarization task
    thread = openai.beta.threads.create(
        messages=[
            {
                "role": "user",
                "content": "Generate a concise summary of the following document",
                "attachments": [
                    {"file_id": messages_file.id, "tools": [{"type": "file_search"}]}
                ],
            }
        ]
    )

    # Run the assistant and fetch the results
    run = openai.beta.threads.runs.create_and_poll(
        thread_id=thread.id, assistant_id=assistant.id
    )

    print(thread)
    print(run)
    
    messages = list(
        openai.beta.threads.messages.list(thread_id=thread.id, run_id=run.id)
    )
    print(messages)
    message_content = messages[0].content[0].text
    annotations = message_content.annotations

    # Annotate the summary content if necessary
    for index, annotation in enumerate(annotations):
        message_content.value = message_content.value.replace(
            annotation.text, f"[{index}]"
        )

    # Insert metadata into the database
    component_metadata = {
        "flow_id": ObjectId(flow_id),
        "name": file.filename,
        "file_id": messages_file.id,
        "assistant_id": assistant.id,
        "vector_store_id": vector_store.id,
        "size": len(file_bytes),
        "type": file_extension,
        "processing_type": "gpt",
        "summary": message_content.value,
    }

    # Store the document in MongoDB
    component_id = component_collection.insert_one(component_metadata).inserted_id

    # Return the component ID with the type
    return {"component_id": str(component_id), "type": file_extension, flow_type: flow_type}

def openai_mindmap_generator(file: UploadFile, flow_id: str, flow_type: str):
    file.file.seek(0)
    file_bytes = file.file.read()
    file_extension = file.filename.split(".")[-1].lower()

    if len(file_bytes) == 0:
        raise ValueError("The uploaded file is actually empty!")

    assistant = openai.beta.assistants.create(
        name="MindMap Builder",
        instructions="Your task is to create the mindmap of the document",
        model="gpt-4o",
        tools=[{"type": "file_search"}],
    )
    vector_store = openai.beta.vector_stores.create(name=f"{file_extension}_mindmap_{flow_id}")

    assistant = openai.beta.assistants.update(
        assistant_id=assistant.id,
        tool_resources={"file_search": {"vector_store_ids": [vector_store.id]}},
    )

    if file_extension == "pdf":
        mime_type = "application/pdf"
        messages_file = openai.files.create(
            file=(file.filename, file_bytes, mime_type), purpose="assistants"
        )
    elif file_extension == "csv":
        mime_type = "application/json"
        df = pd.read_csv(BytesIO(file_bytes))
        json_data = df.to_dict(orient="records")
        json_str = json.dumps(json_data)
        messages_file = openai.files.create(
            file=("data.json", json_str, mime_type), purpose="assistants"
        )
    elif file_extension == "txt":
        mime_type = "text/plain"
        messages_file = openai.files.create(
            file=(file.filename, file_bytes, mime_type), purpose="assistants"
        )
    elif file_extension == "docx":
        mime_type = (
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        )
        messages_file = openai.files.create(
            file=(file.filename, file_bytes, mime_type), purpose="assistants"
        )
    elif file_extension == "html":
        mime_type = "text/html"
        messages_file = openai.files.create(
            file=(file.filename, file_bytes, mime_type), purpose="assistants"
        )
    elif file_extension == "md":
        mime_type = "text/markdown"
        messages_file = openai.files.create(
            file=(file.filename, file_bytes, mime_type), purpose="assistants"
        )
    elif file_extension == "pptx":
        mime_type = (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
        messages_file = openai.files.create(
            file=(file.filename, file_bytes, mime_type), purpose="assistants"
        )
    else:
        raise ValueError("Unsupported file type")
    print("======================================", file_extension)
    thread = openai.beta.threads.create(
        messages=[
            {
                "role": "user",
                "content" : f"""
                You are tasked with generating a JSON mind map that is compatible with React Flow for rendering a flow diagram which should cover all the details and important aspects of the component for which multiple nodes can be required. The mind map should adhere to the following rules:

                1. **Node Types:**
                - There will always be one `dataSource` node, which serves as the root of the flow.
                - There will be `question` node which will be connected to the subsequent `response` node.
                - The `question` node can be connected to data sources or other `response` nodes.
                - There will be `response` for the above question
                
                2. **Node Relationships:**
                - `response` nodes may also connect to each other if it improves the logical flow or visualization.
                - `question` node will always have a `response` node
                - `dataSource` node will always be connected to a question node

                3. **Node Properties:**
                - Each node should have:
                    - `id` (unique identifier of 12 or 24 digit unique uuid or nanoid)
                    - `type` (`dataSource` or `response`)
                    - `position` (coordinates in the form {{ "x": <number>, "y": <number> }} for layout)
                    - `measured` (an object defining width and height):
                        {{
                            "width": <number>,
                            "height": <number>
                        }}
                    - `targetPosition` (position of the target connection, default to `"left"`)
                    - `sourcePosition` (position of the source connection, default to `"right"`)
                    - `selected` (boolean, default to `false`)
                    - `deletable` (boolean, default to `true` for `response` and `false` for `dataSource`)

                4. **Node Data Format:**
                - `dataSource` Node:
                    - `data` contains the following properties:
                        {{
                            "prompt": "<data source description>",
                            "name": "{file_extension}", !!!DOESN"T CHANGES 
                            "content": "<file name or content>",
                            "flow_id": "{flow_id}",
                            "file": "{file.filename}"  // Empty object or file metadata
                        }}
                5. **Question Data Format:**
                - `question` Node:
                    - `data` contains the following properties:
                        {{
                            "question": "<the question asked for the response>",
                            "component_id": "<component reference ID - unique identifier of 12 or 24 digit unique uuid or nanoid>",
                            "component_type" : "{file_extension}",
                        }}
                6. **RESPONSE NODE FORMAT**
                - `response` Node:
                    - `data` contains nested properties:
                        {{
                            "id": "<unique identifier of 12 or 24 digit unique uuid or nanoid>",
                            "type": "response" !!DOESN'T CHANGE,
                            "data": {{
                                "question": "<question text, if applicable>",
                                "summ": "<!!give me a detailed answer for the above question>",
                                "df": [],
                                "graph": "",
                                "flow_id": "{flow_id}",
                                "component_id": "<component reference ID - unique identifier of 12 or 24 digit unique uuid or nanoid>",
                                "component_type": "{file_extension}"
                            }}
                        }}

                7. **Connections:**
                - Connections between nodes should be represented by edges, with the following format:
                    - `id` (unique identifier for the edge)
                    - `source` (ID of the source node)
                    - `target` (ID of the target node)
                    - `type` (optional, defaults to `default`)
                    - 'animated' !!WILL ALWAYS BE TRUE

                8. **Viewport Configuration:**
                - Include a `viewport` object that specifies:
                    - `x` (horizontal position of the viewport)
                    - `y` (vertical position of the viewport)
                    - `zoom` (zoom level for initial rendering)

                ### Additional Considerations:
                - Ensure that the node positions are distributed properly to avoid overlap.
                - Prioritize connecting `response` nodes where it adds logical structure to the flow.

                ### IMPORTANT:
                - **RETURN ONLY THE VALID JSON OBJECT AND NO ADDITIONAL COMMENTS**.
                - Do **not** include any explanations, text, or additional information.
                - Maintain the format with double curly braces `{{` and `}}` as shown in the format.
                """,     

                "attachments": [
                    {"file_id": messages_file.id, "tools": [{"type": "file_search"}]}
                ],
            }
        ]
    )

    # Run the assistant and fetch the results
    run = openai.beta.threads.runs.create_and_poll(
        thread_id=thread.id, assistant_id=assistant.id
    )

    messages = list(
        openai.beta.threads.messages.list(thread_id=thread.id, run_id=run.id)
    )
    message_content = messages[0].content[0].text
    annotations = message_content.annotations

    # Annotate the summary content if necessary
    for index, annotation in enumerate(annotations):
        message_content.value = message_content.value.replace(
            annotation.text, f"[{index}]"
        )

    print(message_content.value)

    response_json = (
            message_content.value.replace("```json", "")
            .replace("```", "")
            .replace("\n", "")
            .strip()
        )

    response_json = json.loads(response_json)

    component_metadata = {
        "flow_id": ObjectId(flow_id),
        "name": file.filename,
        "file_id": messages_file.id,
        "assistant_id": assistant.id,
        "vector_store_id": vector_store.id,
        "size": len(file_bytes),
        "type": file_extension,
        "processing_type": "gpt",
        "mindmap_json": response_json,
    }

    component_id = component_collection.insert_one(component_metadata).inserted_id
    flow = flow_collection.find_one({"_id": ObjectId(flow_id)})
    return {
        "flow_id": flow_id,
        "flow_name": flow["flow_name"],
        "component_id": str(component_id),
        "type": file_extension,
        "mindmap_json": response_json,
        "flow_type": flow_type
    }

    
def one_shot_openai(query, vector_store_id, file_id, assistant_id):
    try:
        template = f"""
        You are an AI assistant tasked with answering the user’s question based on the provided conversation history. Return the results in **JSON format** with the structure below:  

        #### **Response Format:**  
        {{
        "summ": "Your summarized response here...",
        "df": an array of JSON objects,
        "graph": "json_string_representation_of_plotly_graph"
        }}

        ### **Instructions:**
        1. Answer the question using the conversation history.
        2. Extract relevant tabular data into a JSON object compatible with Ag-Grid. If no table exists, return empty JSON object.
        3. If a dataframe is available, generate a relevant **Plotly graph**. Return it as a **valid JSON string** that can be parsed in React.js.
        4. If no graph is possible, return an empty string `""`.
        5. ** The graph's background will be black, so adjust the theme accordingly**.

        NOTE -- "Make sure you need to return only json as response only & please don't add any comments"
        NOTE -- "Make sure you need only need the answer for which context of data is available if not available return empty json as per format"

        **Here is the question:** {query}  

        ### **Example Output:**  
        If the conversation history contains a table and a relevant graph, return:  

        ```json
        {{
        "summ": "Based on the conversation, the key points discussed were...",
        "df": [
            {{
            "column1": "value1",
            "column2": "value2",
            "column3": "value3"
            }},
            {{
            "column1": "value1",
            "column2": "value2",
            "column3": "value3"
            }}
        ],
        "graph": "{{\"data\": [{{\"x\": [\"2024-02-01\", \"2024-02-02\"], \"y\": [100, 150], \"type\": \"line\"}}], \"layout\": {{\"title\": \"Sample Graph\"}}"
        }}
        """

        print(file_id)
        print(assistant_id)
        print(template)
        thread = openai.beta.threads.create(
            messages=[
                {
                    "role": "user",
                    "content": template,
                    "attachments": [
                        {"file_id": file_id, "tools": [{"type": "file_search"}]}
                    ],
                }
            ]
        )
        run = openai.beta.threads.runs.create_and_poll(
            thread_id=thread.id, assistant_id=assistant_id
        )

        messages = list(
            openai.beta.threads.messages.list(thread_id=thread.id, run_id=run.id)
        )
        message_content = messages[0].content[0].text
        print(message_content)
        annotations = message_content.annotations
        for index, annotation in enumerate(annotations):
            message_content.value = message_content.value.replace(
                annotation.text, f"[{index}]"
            )
            message_content.value = (
                message_content.value.replace("```json", "").replace("```", "").strip()
            )
            message_content.value = message_content.value.replace("\n", "")
        response = message_content.value.replace("```json", "").replace("```", "").strip().replace("\n", "") 
        print(response)
        return response
    except Exception as e:
        print(e.with_traceback())
        raise HTTPException(status_code=500)


def get_page_len(file: UploadFile):
    f_bytes = file.file.read()
    stream = BytesIO(f_bytes)
    reader = pdfium.PdfDocument(f_bytes)
    stream.seek(0)
    file.file.seek(0)
    if len(reader) > 100:
        return True
    else:
        return False


@app.post("/component-create-pdf")
def create_pdf_component(
    file: UploadFile, flow_id: str = Form(...), processing_type: str = Form(...)
):
    flow = flow_collection.find_one({"_id": ObjectId(flow_id)})
    if file.filename.endswith(".pdf"):
        print(get_page_len(file))
        check_page_length = get_page_len(file)
        if processing_type == "gpt" and not check_page_length and flow["flow_type"] == 'manual':
            return get_summary_from_openai(file, flow_id=flow_id, flow_type='manual')
        elif processing_type == "aws" and flow["flow_type"] == 'manual':
            return use_aws_textract(file, flow_id=flow_id, flow_type='manual')
        elif processing_type == "custom" and flow["flow_type"] == 'manual':
            return camelot_pdf_processing(flow_id, file, 'manual')
        elif processing_type == "gpt" and not check_page_length and flow["flow_type"] == 'automatic':
            return openai_mindmap_generator(file, flow_id=flow_id, flow_type='automatic')
        elif processing_type == "aws" and flow["flow_type"] == 'automatic':
            return use_aws_textract(file, flow_id=flow_id, flow_type='automatic')
        elif processing_type == "custom" and flow["flow_type"] == 'automatic':
            return camelot_pdf_processing(flow_id, file, 'automatic')
        else:
            traceback.print_exc()
            return HTTPException(status_code=404, detail="Exceeded Page limit for GPT.")
    else:
        traceback.print_exc()
        raise HTTPException(status_code=400, detail="Only PDF files are allowed.")

@app.post("/component-create-img")
async def create_img_component(flow_id: str = Form(...), file: UploadFile = File(...)):
    try:
        flow = flow_collection.find_one({"_id": ObjectId(flow_id)})
        MAX_IMAGE_SIZE_MB = 16
        ALLOWED_MIME_TYPES = {
            "image/jpeg",
            "image/png",
            "image/webp",
            "image/heic",
            "image/heif",
        }

        if file.content_type not in ALLOWED_MIME_TYPES:
            raise HTTPException(
                status_code=400, detail=f"Unsupported file type: {file.content_type}"
            )

        contents = await file.read()

        size_in_mb = len(contents) / (1024 * 1024)
        if size_in_mb > MAX_IMAGE_SIZE_MB:
            raise HTTPException(status_code=400, detail="Image exceeds 16MB size limit")

        image_base64 = base64.b64encode(contents).decode("utf-8")
        
        if flow["flow_type"] == 'manual':

            component_metadata = {
                "flow_id": ObjectId(flow_id),
                "name": file.filename,
                "mime_type": file.content_type,
                "type": "image",
                "base64_image": image_base64,
                "processing_type": "claude" if use_bedrock_only else "gemini",
            }

            component_id = component_collection.insert_one(component_metadata).inserted_id

            return {
                "message": "Image component created successfully",
                "component_id": str(component_id),
                "type": "image",
            }
        
        else:
            
            template = f"""You are tasked with generating a JSON mind map for given image file and that should be compatible with React Flow for rendering a flow diagram which should cover all the details and important aspects of the component for which multiple nodes can be required. The mind map should adhere to the following rules:

                1. **Node Types:**
                - There will always be one `dataSource` node, which serves as the root of the flow.
                - There will be `question` node which will be connected to the subsequent `response` node.
                - The `question` node can be connected to data sources or other `response` nodes.
                - There will be `response` for the above question
                
                2. **Node Relationships:**
                - `response` nodes may also connect to each other if it improves the logical flow or visualization.
                - `question` node will always have a `response` node
                - `dataSource` node will always be connected to a question node

                3. **Node Properties:**
                - Each node should have:
                    - `id` (unique identifier of 12 or 24 digit unique uuid or nanoid)
                    - `type` (`dataSource` or `response`)
                    - `position` (coordinates in the form {{ "x": <number>, "y": <number> }} for layout)
                    - `measured` (an object defining width and height):
                        {{
                            "width": <number>,
                            "height": <number>
                        }}
                    - `targetPosition` (position of the target connection, default to `"left"`)
                    - `sourcePosition` (position of the source connection, default to `"right"`)
                    - `selected` (boolean, default to `false`)
                    - `deletable` (boolean, default to `true` for `response` and `false` for `dataSource`)

                4. **Node Data Format:**
                - `dataSource` Node:
                    - `data` contains the following properties:
                        {{
                            "prompt": "<data source description>",
                            "name": "image", !!!DOESN"T CHANGES 
                            "content": "<file name or content>",
                            "flow_id": "{flow_id}",
                            "file": "{file.filename}"  // Empty object or file metadata
                        }}
                5. **Question Data Format:**
                - `question` Node:
                    - `data` contains the following properties:
                        {{
                            "question": "<the question asked for the response>",
                            "component_id": "<component reference ID - unique identifier of 12 or 24 digit unique uuid or nanoid>",
                            "component_type" : "image",
                        }}
                6. **RESPONSE NODE FORMAT**
                - `response` Node:
                    - `data` contains nested properties:
                        {{
                            "id": "<unique identifier of 12 or 24 digit unique uuid or nanoid>",
                            "type": "response" !!DOESN'T CHANGE,
                            "data": {{
                                "question": "<question text, if applicable>",
                                "summ": "<!!give me a detailed answer for the above question>",
                                "df": [],
                                "graph": "",
                                "flow_id": "{flow_id}",
                                "component_id": "<component reference ID - unique identifier of 12 or 24 digit unique uuid or nanoid>",
                                "component_type": "image"
                            }}
                        }}

                7. **Connections:**
                - Connections between nodes should be represented by edges, with the following format:
                    - `id` (unique identifier for the edge)
                    - `source` (ID of the source node)
                    - `target` (ID of the target node)
                    - `type` (optional, defaults to `default`)
                    - 'animated' !!WILL ALWAYS BE TRUE

                8. **Viewport Configuration:**
                - Include a `viewport` object that specifies:
                    - `x` (horizontal position of the viewport)
                    - `y` (vertical position of the viewport)
                    - `zoom` (zoom level for initial rendering)

                ### Additional Considerations:
                - Ensure that the node positions are distributed properly to avoid overlap.
                - Prioritize connecting `response` nodes where it adds logical structure to the flow.

                ### IMPORTANT:
                - **RETURN ONLY THE VALID JSON OBJECT AND NO ADDITIONAL COMMENTS**.
                - Do **not** include any explanations, text, or additional information.
                - Maintain the format with double curly braces `{{` and `}}` as shown in the format.
                """
            
            if use_bedrock_only:
                # Use Claude for image analysis
                image_base64 = base64.b64encode(contents).decode("utf-8")
                response_text = analyze_image_with_claude(image_base64, template, file.content_type)
            else:
                # Use Gemini for image analysis
                image_part = {"mime_type": file.content_type, "data": contents}
                response = model.generate_content(contents=[template, image_part])
                response_text = response.text

            response_json = response_text.replace("```json", "").replace("```", "").replace("\n", "").strip()
            response_json = json.loads(response_json)
        
        print(response_json)
        
        component_metadata = {
            "flow_id": flow_id,
            "name": file.filename,
            "type": "image",
            "processing_type": "claude" if use_bedrock_only else "gemini",
            "mindmap_json": response_json,
        }

        component_id = component_collection.insert_one(component_metadata).inserted_id

        return {
            "flow_id" : ObjectId(flow_id),
            "flow_name": flow["flow_name"],
            "component_id": str(component_id),
            "type": "image",
            "mindmap_json": response_json,
            "flow_type": "automatic"
        }  

    except Exception as e:
        print(f"Error in /component-create-img endpoint: {e}")
        raise HTTPException(status_code=500, detail="Internal Server Error")


@app.post("/component-create-audio")
async def create_audio_component(
    flow_id: str = Form(...), file: UploadFile = File(...)
):
    try:
        flow = flow_collection.find_one({"_id": ObjectId(flow_id)})
        MAX_AUDIO_SIZE_MB = 16
        ALLOWED_MIME_TYPES = {
            "audio/wav",
            "audio/mp3",
            "audio/aiff",
            "audio/aac",
            "audio/ogg",
            "audio/flac",
            "audio/mpeg",
        }

        if file.content_type not in ALLOWED_MIME_TYPES:
            raise HTTPException(
                status_code=400,
                detail=f"Unsupported audio file type: {file.content_type}",
            )

        contents = await file.read()
        size_in_mb = len(contents) / (1024 * 1024)

        if size_in_mb > MAX_AUDIO_SIZE_MB:
            raise HTTPException(status_code=400, detail="Audio exceeds 16MB size limit")

        audio_base64 = base64.b64encode(contents).decode("utf-8")
        
        if flow["flow_type"] == 'manual':

            component_metadata = {
                "flow_id": ObjectId(flow_id),
                "name": file.filename,
                "mime_type": file.content_type,
                "type": "audio",
                "base64_audio": audio_base64,
                "processing_type": "pegasus" if use_bedrock_only else "gemini",
            }

            component_id = component_collection.insert_one(component_metadata).inserted_id

            return {
                "message": "Audio component created successfully",
                "component_id": str(component_id),
                "type": "audio",
            }
            
        else:
            
            template = f"""You are tasked with generating a JSON mind map for given audio file and that should be compatible with React Flow for rendering a flow diagram which should cover all the details and important aspects of the component for which multiple nodes can be required. The mind map should adhere to the following rules:

                1. **Node Types:**
                - There will always be one `dataSource` node, which serves as the root of the flow.
                - There will be `question` node which will be connected to the subsequent `response` node.
                - The `question` node can be connected to data sources or other `response` nodes.
                - There will be `response` for the above question
                
                2. **Node Relationships:**
                - `response` nodes may also connect to each other if it improves the logical flow or visualization.
                - `question` node will always have a `response` node
                - `dataSource` node will always be connected to a question node

                3. **Node Properties:**
                - Each node should have:
                    - `id` (unique identifier of 12 or 24 digit unique uuid or nanoid)
                    - `type` (`dataSource` or `response`)
                    - `position` (coordinates in the form {{ "x": <number>, "y": <number> }} for layout)
                    - `measured` (an object defining width and height):
                        {{
                            "width": <number>,
                            "height": <number>
                        }}
                    - `targetPosition` (position of the target connection, default to `"left"`)
                    - `sourcePosition` (position of the source connection, default to `"right"`)
                    - `selected` (boolean, default to `false`)
                    - `deletable` (boolean, default to `true` for `response` and `false` for `dataSource`)

                4. **Node Data Format:**
                - `dataSource` Node:
                    - `data` contains the following properties:
                        {{
                            "prompt": "<data source description>",
                            "name": "audio", !!!DOESN"T CHANGES 
                            "content": "<file name or content>",
                            "flow_id": "{flow_id}",
                            "file": "{file.filename}"  // Empty object or file metadata
                        }}
                5. **Question Data Format:**
                - `question` Node:
                    - `data` contains the following properties:
                        {{
                            "question": "<the question asked for the response>",
                            "component_id": "<component reference ID - unique identifier of 12 or 24 digit unique uuid or nanoid>",
                            "component_type" : "audio",
                        }}
                6. **RESPONSE NODE FORMAT**
                - `response` Node:
                    - `data` contains nested properties:
                        {{
                            "id": "<unique identifier of 12 or 24 digit unique uuid or nanoid>",
                            "type": "response" !!DOESN'T CHANGE,
                            "data": {{
                                "question": "<question text, if applicable>",
                                "summ": "<!!give me a detailed answer for the above question>",
                                "df": [],
                                "graph": "",
                                "flow_id": "{flow_id}",
                                "component_id": "<component reference ID - unique identifier of 12 or 24 digit unique uuid or nanoid>",
                                "component_type": "audio"
                            }}
                        }}

                7. **Connections:**
                - Connections between nodes should be represented by edges, with the following format:
                    - `id` (unique identifier for the edge)
                    - `source` (ID of the source node)
                    - `target` (ID of the target node)
                    - `type` (optional, defaults to `default`)
                    - 'animated' !!WILL ALWAYS BE TRUE

                8. **Viewport Configuration:**
                - Include a `viewport` object that specifies:
                    - `x` (horizontal position of the viewport)
                    - `y` (vertical position of the viewport)
                    - `zoom` (zoom level for initial rendering)

                ### Additional Considerations:
                - Ensure that the node positions are distributed properly to avoid overlap.
                - Prioritize connecting `response` nodes where it adds logical structure to the flow.

                ### IMPORTANT:
                - **RETURN ONLY THE VALID JSON OBJECT AND NO ADDITIONAL COMMENTS**.
                - Do **not** include any explanations, text, or additional information.
                - Maintain the format with double curly braces `{{` and `}}` as shown in the format.
                """

            if use_bedrock_only:
                # Upload audio to S3 and use Pegasus for analysis
                unique_id = str(uuid4())
                s3_key = f"uploads/{flow_id}/audio/{unique_id}_{file.filename}"
                upload_to_s3(contents, bucket_name, s3_key)
                s3_uri = f"s3://{bucket_name}/{s3_key}"
                
                response_text = analyze_content_with_pegasus(s3_uri, template)
            else:
                # Use Gemini for audio analysis
                audio_part = {"mime_type": file.content_type, "data": contents}
                response = model.generate_content(contents=[template, audio_part])
                response_text = response.text

            response_json = response_text.replace("```json", "").replace("```", "").replace("\n", "").strip()
            response_json = json.loads(response_json) 
        print(response_json)
        
        component_metadata = {
            "flow_id": ObjectId(flow_id),
            "name": file.filename,
            "type": "audio",
            "processing_type": "pegasus" if use_bedrock_only else "gemini",
            "mindmap_json": response_json,
        }

        component_id = component_collection.insert_one(component_metadata).inserted_id

        return {
            "flow_id" : flow_id,
            "flow_name": flow["flow_name"],
            "component_id": str(component_id),
            "type": "audio",
            "mindmap_json": response_json,
            "flow_type": "automatic"
        }

    except Exception as e:
        print(f"Error in /component-create-audio endpoint: {e}")
        raise HTTPException(status_code=500, detail="Internal Server Error")


@app.post("/component-create-youtube")
def create_youtube_component(
    flow_id: str = Form(...), youtube_url: str = Form(...)
):
    try:
        flow = flow_collection.find_one({"_id": ObjectId(flow_id)})
        
        print(youtube_url)
        
        if flow["flow_type"] == 'manual':
            component_metadata = {
                "flow_id": ObjectId(flow_id),
                "youtube_url": youtube_url,
                "type": "youtube",
                "processing_type": "pegasus" if use_bedrock_only else "gemini",
            }

            component_id = component_collection.insert_one(component_metadata).inserted_id

            return {
                "message": "Youtube component created successfully",
                "component_id": str(component_id),
                "type": "youtube",
            }
            
        else:
            
            template = f"""
                You are tasked with generating a JSON mind map for give youtube URL and should be compatible with React Flow for rendering a flow diagram which should cover all the details and important aspects of the component for which multiple nodes can be required. The mind map should adhere to the following rules:

                1. **Node Types:**
                - There will always be one `dataSource` node, which serves as the root of the flow.
                - There will be `question` node which will be connected to the subsequent `response` node.
                - The `question` node can be connected to data sources or other `response` nodes.
                - There will be `response` for the above question
                
                2. **Node Relationships:**
                - `response` nodes may also connect to each other if it improves the logical flow or visualization.
                - `question` node will always have a `response` node
                - `dataSource` node will always be connected to a question node

                3. **Node Properties:**
                - Each node should have:
                    - `id` (unique identifier of 12 or 24 digit unique uuid or nanoid)
                    - `type` (`dataSource` or `response`)
                    - `position` (coordinates in the form {{ "x": <number>, "y": <number> }} for layout)
                    - `measured` (an object defining width and height):
                        {{
                            "width": <number>,
                            "height": <number>
                        }}
                    - `targetPosition` (position of the target connection, default to `"left"`)
                    - `sourcePosition` (position of the source connection, default to `"right"`)
                    - `selected` (boolean, default to `false`)
                    - `deletable` (boolean, default to `true` for `response` and `false` for `dataSource`)

                4. **Node Data Format:**
                - `dataSource` Node:
                    - `data` contains the following properties:
                        {{
                            "prompt": "<data source description>",
                            "name": "youtube", !!!DOESN"T CHANGES 
                            "content": "<file name or content>",
                            "flow_id": "{flow_id}",
                            "file": "{youtube_url}"  // Empty object or file metadata
                        }}
                5. **Question Data Format:**
                - `question` Node:
                    - `data` contains the following properties:
                        {{
                            "question": "<the question asked for the response>",
                            "component_id": "<component reference ID - unique identifier of 12 or 24 digit unique uuid or nanoid>",
                            "component_type" : "youtube",
                        }}
                6. **RESPONSE NODE FORMAT**
                - `response` Node:
                    - `data` contains nested properties:
                        {{
                            "id": "<unique identifier of 12 or 24 digit unique uuid or nanoid>",
                            "type": "response" !!DOESN'T CHANGE,
                            "data": {{
                                "question": "<question text, if applicable>",
                                "summ": "<!!give me a detailed answer for the above question>",
                                "df": [],
                                "graph": "",
                                "flow_id": "{flow_id}",
                                "component_id": "<component reference ID - unique identifier of 12 or 24 digit unique uuid or nanoid>",
                                "component_type": "youtube"
                            }}
                        }}

                7. **Connections:**
                - Connections between nodes should be represented by edges, with the following format:
                    - `id` (unique identifier for the edge)
                    - `source` (ID of the source node)
                    - `target` (ID of the target node)
                    - `type` (optional, defaults to `default`)
                    - 'animated' !!WILL ALWAYS BE TRUE

                8. **Viewport Configuration:**
                - Include a `viewport` object that specifies:
                    - `x` (horizontal position of the viewport)
                    - `y` (vertical position of the viewport)
                    - `zoom` (zoom level for initial rendering)

                ### Additional Considerations:
                - Ensure that the node positions are distributed properly to avoid overlap.
                - Prioritize connecting `response` nodes where it adds logical structure to the flow.

                ### IMPORTANT:
                - **RETURN ONLY THE VALID JSON OBJECT AND NO ADDITIONAL COMMENTS**.
                - Do **not** include any explanations, text, or additional information.
                - Maintain the format with double curly braces `{{` and `}}` as shown in the format.
                """

            if use_bedrock_only:
                # Use Pegasus for YouTube analysis
                # Note: For YouTube URLs, we need to construct the S3 URI differently
                # or use the direct URL if Pegasus supports it
                response_text = analyze_content_with_pegasus(youtube_url, template)
            else:
                # Use Gemini for YouTube analysis
                mime_type = "video/*"
                response = model_vertexai.generate_content(
                    contents=[template, Part.from_uri(youtube_url, mime_type)]
                )
                response_text = response.text

            response_json = response_text.replace("```json", "").replace("```", "").replace("\n", "").strip()
            response_json = json.loads(response_json)
        
        print(response_json)
        
        component_metadata = {
            "flow_id": ObjectId(flow_id),
            "youtube_url": youtube_url,
            "type": "youtube",
            "processing_type": "pegasus" if use_bedrock_only else "gemini",
            "mindmap_json": response_json,
        }

        component_id = component_collection.insert_one(component_metadata).inserted_id

        return {
            "flow_id" : flow_id,
            "flow_name": flow["flow_name"],
            "component_id": str(component_id),
            "type": "youtube",
            "mindmap_json": response_json,
            "flow_type": "automatic"
        }
            

    except Exception as e:
        traceback.print_exc()
        print(f"Error in /component-create-youtube endpoint: {e}")
        raise HTTPException(status_code=500, detail="Internal Server Error")


@app.post("/component-create-video")
async def create_video_component(
    flow_id: str = Form(...), file: UploadFile = File(...)
):
    try:
        flow = flow_collection.find_one({"_id": ObjectId(flow_id)})
        MAX_VIDEO_SIZE_MB = 16
        ALLOWED_MIME_TYPES = {
            "video/x-flv",
            "video/quicktime",
            "video/mpeg",
            "video/mpegs",
            "video/mpgs",
            "video/mpg",
            "video/mp4",
            "video/webm",
            "video/wmv",
            "video/3gpp",
        }

        if file.content_type not in ALLOWED_MIME_TYPES:
            raise HTTPException(
                status_code=400,
                detail=f"Unsupported video file type: {file.content_type}",
            )

        contents = await file.read()
        size_in_mb = len(contents) / (1024 * 1024)

        if size_in_mb > MAX_VIDEO_SIZE_MB:
            raise HTTPException(status_code=400, detail="Video exceeds 16MB size limit")

        unique_id = str(uuid4())
        s3_key = f"uploads/{flow_id}/videos/{unique_id}_{file.filename}"

        upload_to_s3(contents, bucket_name, s3_key)

        video_url = f"https://{bucket_name}.s3.amazonaws.com/{s3_key}"
        
        if flow["flow_type"] == 'manual':

            component_metadata = {
                "flow_id": ObjectId(flow_id),
                "name": file.filename,
                "mime_type": file.content_type,
                "type": "video",
                "video_url": video_url,
                "processing_type": "pegasus" if use_bedrock_only else "gemini",
            }

            component_id = component_collection.insert_one(component_metadata).inserted_id

            return {
                "message": "Video component created successfully",
                "component_id": str(component_id),
                "type": "video",
            }
            
        else:
            
            template = f"""
                You are tasked with generating a JSON mind map for give video and should be compatible with React Flow for rendering a flow diagram which should cover all the details and important aspects of the component for which multiple nodes can be required. The mind map should adhere to the following rules:

                1. **Node Types:**
                - There will always be one `dataSource` node, which serves as the root of the flow.
                - There will be `question` node which will be connected to the subsequent `response` node.
                - The `question` node can be connected to data sources or other `response` nodes.
                - There will be `response` for the above question
                
                2. **Node Relationships:**
                - `response` nodes may also connect to each other if it improves the logical flow or visualization.
                - `question` node will always have a `response` node
                - `dataSource` node will always be connected to a question node

                3. **Node Properties:**
                - Each node should have:
                    - `id` (unique identifier of 12 or 24 digit unique uuid or nanoid)
                    - `type` (`dataSource` or `response`)
                    - `position` (coordinates in the form {{ "x": <number>, "y": <number> }} for layout)
                    - `measured` (an object defining width and height):
                        {{
                            "width": <number>,
                            "height": <number>
                        }}
                    - `targetPosition` (position of the target connection, default to `"left"`)
                    - `sourcePosition` (position of the source connection, default to `"right"`)
                    - `selected` (boolean, default to `false`)
                    - `deletable` (boolean, default to `true` for `response` and `false` for `dataSource`)

                4. **Node Data Format:**
                - `dataSource` Node:
                    - `data` contains the following properties:
                        {{
                            "prompt": "<data source description>",
                            "name": "{file.filename}", !!!DOESN"T CHANGES 
                            "content": "<file name or content>",
                            "flow_id": "{flow_id}",
                            "file": "{file.filename}"  // Empty object or file metadata
                        }}
                5. **Question Data Format:**
                - `question` Node:
                    - `data` contains the following properties:
                        {{
                            "question": "<the question asked for the response>",
                            "component_id": "<component reference ID - unique identifier of 12 or 24 digit unique uuid or nanoid>",
                            "component_type" : "video",
                        }}
                6. **RESPONSE NODE FORMAT**
                - `response` Node:
                    - `data` contains nested properties:
                        {{
                            "id": "<unique identifier of 12 or 24 digit unique uuid or nanoid>",
                            "type": "response" !!DOESN'T CHANGE,
                            "data": {{
                                "question": "<question text, if applicable>",
                                "summ": "<!!give me a detailed answer for the above question>",
                                "df": [],
                                "graph": "",
                                "flow_id": "{flow_id}",
                                "component_id": "<component reference ID - unique identifier of 12 or 24 digit unique uuid or nanoid>",
                                "component_type": "video"
                            }}
                        }}

                7. **Connections:**
                - Connections between nodes should be represented by edges, with the following format:
                    - `id` (unique identifier for the edge)
                    - `source` (ID of the source node)
                    - `target` (ID of the target node)
                    - `type` (optional, defaults to `default`)
                    - 'animated' !!WILL ALWAYS BE TRUE

                8. **Viewport Configuration:**
                - Include a `viewport` object that specifies:
                    - `x` (horizontal position of the viewport)
                    - `y` (vertical position of the viewport)
                    - `zoom` (zoom level for initial rendering)

                ### Additional Considerations:
                - Ensure that the node positions are distributed properly to avoid overlap.
                - Prioritize connecting `response` nodes where it adds logical structure to the flow.

                ### IMPORTANT:
                - **RETURN ONLY THE VALID JSON OBJECT AND NO ADDITIONAL COMMENTS**.
                - Do **not** include any explanations, text, or additional information.
                - Maintain the format with double curly braces `{{` and `}}` as shown in the format.
                """

            if use_bedrock_only:
                # Use Pegasus for video analysis
                s3_uri = f"s3://{bucket_name}/{s3_key}"
                response_text = analyze_content_with_pegasus(s3_uri, template)
            else:
                # Use Gemini for video analysis
                mime_type = "video/*"
                response = model_vertexai.generate_content(
                    contents=[template, Part.from_uri(video_url, mime_type)]
                )
                response_text = response.text

            response_json = response_text.replace("```json", "").replace("```", "").replace("\n", "").strip()
            response_json = json.loads(response_json)
        
        print(response_json)
                
        component_metadata = {
            "flow_id": ObjectId(flow_id),
            "video_url": video_url,
            "type": "video",
            "processing_type": "pegasus" if use_bedrock_only else "gemini",
            "mindmap_json": response_json,
        }

        component_id = component_collection.insert_one(component_metadata).inserted_id

        return {
            "flow_id" : flow_id,
            "flow_name": flow["flow_name"],
            "component_id": str(component_id),
            "type": "video",
            "mindmap_json": response_json,
            "flow_type": "automatic"
        }
            
    except Exception as e:
        print(f"Error in /component-create-video endpoint: {e}")
        raise HTTPException(status_code=500, detail="Internal Server Error")


@app.post("/component-create-txt")
def create_txt_component(file: UploadFile, flow_id: str = Form(...)):
    flow = flow_collection.find_one({"_id": ObjectId(flow_id)})
    if file.filename.endswith(".txt"):
        check_page_length = is_within_gpt4o_token_limit(file)
        if check_page_length and flow["flow_type"] == 'manual':
            return get_summary_from_openai(file, flow_id=flow_id, flow_type=flow["flow_type"])
        elif check_page_length and flow["flow_type"] == 'automatic':
            return openai_mindmap_generator(file, flow_id=flow_id, flow_type=flow["flow_type"])
        else:
            traceback.print_exc()
            return HTTPException(status_code=404, detail="Exceeded Page limit for GPT.")
    else:
        traceback.print_exc()
        raise HTTPException(status_code=400, detail="Only TXT files are allowed.")


@app.post("/component-create-md")
def create_md_component(file: UploadFile, flow_id: str = Form(...)):
    flow = flow_collection.find_one({"_id": ObjectId(flow_id)})
    if file.filename.endswith(".md"):
        check_page_length = is_within_gpt4o_token_limit(file)
        if check_page_length and flow["flow_type"] == 'manual':
            return get_summary_from_openai(file, flow_id=flow_id, flow_type=flow["flow_type"])
        elif check_page_length and flow["flow_type"] == 'automatic':
            return openai_mindmap_generator(file, flow_id=flow_id, flow_type=flow["flow_type"])
        else:
            traceback.print_exc()
            return HTTPException(status_code=404, detail="Exceeded Page limit for GPT.")
    else:
        traceback.print_exc()
        raise HTTPException(status_code=400, detail="Only MarkDown files are allowed.")


@app.post("/component-create-pptx")
def create_pptx_component(file: UploadFile, flow_id: str = Form(...)):
    flow = flow_collection.find_one({"_id": ObjectId(flow_id)})
    if file.filename.endswith(".pptx"):
        check_page_length = is_within_gpt4o_token_limit(file)
        if check_page_length and flow["flow_type"] == 'manual':
            return get_summary_from_openai(file, flow_id=flow_id, flow_type=flow["flow_type"])
        elif check_page_length and flow["flow_type"] == 'automatic':
            return openai_mindmap_generator(file, flow_id=flow_id, flow_type=flow["flow_type"])
        else:
            traceback.print_exc()
            return HTTPException(status_code=404, detail="Exceeded Page limit for GPT.")
    else:
        traceback.print_exc()
        raise HTTPException(status_code=400, detail="Only PPTX files are allowed.")


@app.post("/component-create-html")
def create_html_component(file: UploadFile, flow_id: str = Form(...)):
    flow = flow_collection.find_one({"_id": ObjectId(flow_id)})
    if file.filename.endswith(".html"):
        check_page_length = is_within_gpt4o_token_limit(file)
        if check_page_length and flow["flow_type"] == 'manual':
            return get_summary_from_openai(file, flow_id=flow_id, flow_type=flow["flow_type"])
        elif check_page_length and flow["flow_type"] == 'automatic':
            return openai_mindmap_generator(file, flow_id=flow_id, flow_type=flow["flow_type"])
        else:
            traceback.print_exc()
            return HTTPException(status_code=404, detail="Exceeded Page limit for GPT.")
    else:
        traceback.print_exc()
        raise HTTPException(status_code=400, detail="Only HTML files are allowed.")


@app.post("/component-create-docx")
def create_docx_component(file: UploadFile, flow_id: str = Form(...)):
    flow = flow_collection.find_one({"_id": ObjectId(flow_id)})
    if file.filename.endswith(".docx"):
        check_page_length = is_within_gpt4o_token_limit(file)
        if check_page_length and flow["flow_type"] == 'manual':
            return get_summary_from_openai(file, flow_id=flow_id, flow_type=flow["flow_type"])
        elif check_page_length and flow["flow_type"] == 'automatic':
            return openai_mindmap_generator(file, flow_id=flow_id, flow_type=flow["flow_type"])
        else:
            traceback.print_exc()
            return HTTPException(status_code=404, detail="Exceeded Page limit for GPT.")
    else:
        traceback.print_exc()
        raise HTTPException(status_code=400, detail="Only DOCX files are allowed.")


@app.post("/component-create-csv")
def create_csv_component(
    file: UploadFile = File(...), flow_id: str = Form(...), header_row: int = Form(...)
):
    try:
        flow = flow_collection.find_one({"_id": ObjectId(flow_id)})
        if flow["flow_type"] != 'manual':
            raise HTTPException(status_code=400, detail="Only Manual Mindmap is supported for CSV.")
        if not file.filename.endswith(".csv"):
            raise HTTPException(status_code=400, detail="Only CSV files are allowed.")

        file_bytes = file.file.read()
        file_hash = calculate_file_hash(file_bytes)

        existing_component = component_collection.find_one(
            {"file_hash": file_hash, "flow_id": ObjectId(flow_id)}
        )
        if existing_component:
            raise HTTPException(
                status_code=400, detail="File already exists in the system."
            )

        unique_table_name = f"tbl_{uuid4().hex[:8]}"

        file_name = file.filename
        folder = f"uploads/{flow_id}/"
        s3_key = folder + file_name
        upload_to_s3(file_bytes, bucket_name, s3_key)
        print("uploaded")

        sql_con = sqlite3.connect("csv_data.db")
        buffer = BytesIO(file_bytes)
        df = pd.read_csv(
            buffer, skiprows=header_row, encoding="utf-8", encoding_errors="ignore"
        )
        print(df)
        buffer.close()
        file.file.close()
        print("CSV into SQLite")

        df.to_sql(name=unique_table_name, con=sql_con, if_exists="replace", index=False)
        sql_con.close()
        csvBot.connect_to_sqlite("csv_data.db")

        df_ddl = csvBot.run_sql(
            f"SELECT sql FROM sqlite_master WHERE type = 'table' AND name = '{unique_table_name}'"
        )

        for ddl in df_ddl["sql"].to_list():
            csvBot.train(ddl=ddl)

        training_data = csvBot.get_training_data()
        print(training_data)

        component_metadata = {
            "flow_id": ObjectId(flow_id),
            "name": file.filename,
            "table_name": unique_table_name,
            "file_hash": file_hash,
            "size": len(file_bytes),
            "type": "csv",
            "s3_path": s3_key,
            "created_at": datetime.datetime.utcnow(),
        }
        component_id = component_collection.insert_one(component_metadata).inserted_id
        return {
            "component_id": str(component_id),
            "type": "csv",
            "message": "Component created successfully",
        }

    except Exception as e:
        print(f"Error in /component-create-csv endpoint: {e}")
        raise HTTPException(status_code=500, detail="Internal Server Error")


@app.post("/component-create-crawl")
async def create_web_crawler(
    flow_id: str = Form(...),
    web_url: str = Form(...),
):
    try:
        flow = flow_collection.find_one({"_id": ObjectId(flow_id)})
        flow_type = flow["flow_type"]
        unique_id = str(uuid4())

        async with AsyncWebCrawler() as crawler:
            result = await crawler.arun(url=web_url)

        response = result.markdown

        file_bytes = response.encode("utf-8")

        mime_type = "text/markdown"
        
        if flow_type == "manual":
            assistant = openai.beta.assistants.create(
                name="Summarize agent",
                instructions="Your task is to only summarize the document",
                model="gpt-4o",
                tools=[{"type": "file_search"}],
            )
            vector_store = openai.beta.vector_stores.create(name=f"web_{flow_id}")

            assistant = openai.beta.assistants.update(
                assistant_id=assistant.id,
                tool_resources={"file_search": {"vector_store_ids": [vector_store.id]}},
            )

            
            messages_file = openai.files.create(
                file=(f"website_{unique_id}.md", file_bytes, mime_type),
                purpose="assistants",
            )

            thread = openai.beta.threads.create(
                messages=[
                {
                    "role": "user",
                    "content": "Generate a concise summary of the following document",
                    "attachments": [
                        {
                            "file_id": messages_file.id,
                            "tools": [{"type": "file_search"}],
                        }
                    ],
                }
                ]
            )

            run = openai.beta.threads.runs.create_and_poll(
                thread_id=thread.id, assistant_id=assistant.id
            )

            messages = list(
                openai.beta.threads.messages.list(thread_id=thread.id, run_id=run.id)
            )
            
            print(messages)
            message_content = messages[0].content[0].text
            annotations = message_content.annotations

            for index, annotation in enumerate(annotations):
                message_content.value = message_content.value.replace(annotation.text, f"[{index}]")

            component_metadata = {
                "flow_id": ObjectId(flow_id),
                "file_id": messages_file.id,
                "assistant_id": assistant.id,
                "vector_store_id": vector_store.id,
                "size": len(file_bytes),
                "type": "web",
                "web_url": web_url,
                "processing_type": "gpt",
                "summary": message_content.value,
            }

            component_id = component_collection.insert_one(component_metadata).inserted_id

            return {
                "component_id": str(component_id),
                "type": "web",
                "message": "Component created successfully",
            }
            
        else:
            assistant = openai.beta.assistants.create(
            name="MindMap Builder",
            instructions="Your task is to create the mindmap of the document",
            model="gpt-4o",
            tools=[{"type": "file_search"}],
            )
            vector_store = openai.beta.vector_stores.create(name=f"web_mindmap_{flow_id}")

            assistant = openai.beta.assistants.update(
            assistant_id=assistant.id,
            tool_resources={"file_search": {"vector_store_ids": [vector_store.id]}},
            )

            messages_file = openai.files.create(
                file=(f"website_mindmap_{unique_id}.md", file_bytes, mime_type), purpose="assistants"
            )

            thread = openai.beta.threads.create(
            
            messages=[
            {
                "role": "user",
                "content" : f"""
                You are tasked with generating a JSON mind map that is compatible with React Flow for rendering a flow diagram which should cover all the details and important aspects of the component for which multiple nodes can be required. The mind map should adhere to the following rules:

                1. **Node Types:**
                - There will always be one `dataSource` node, which serves as the root of the flow.
                - There will be `question` node which will be connected to the subsequent `response` node.
                - The `question` node can be connected to data sources or other `response` nodes.
                - There will be `response` for the above question
                
                2. **Node Relationships:**
                - `response` nodes may also connect to each other if it improves the logical flow or visualization.
                - `question` node will always have a `response` node
                - `dataSource` node will always be connected to a question node

                3. **Node Properties:**
                - Each node should have:
                    - `id` (unique identifier of 12 or 24 digit unique uuid or nanoid)
                    - `type` (`dataSource` or `response`)
                    - `position` (coordinates in the form {{ "x": <number>, "y": <number> }} for layout)
                    - `measured` (an object defining width and height):
                        {{
                            "width": <number>,
                            "height": <number>
                        }}
                    - `targetPosition` (position of the target connection, default to `"left"`)
                    - `sourcePosition` (position of the source connection, default to `"right"`)
                    - `selected` (boolean, default to `false`)
                    - `deletable` (boolean, default to `true` for `response` and `false` for `dataSource`)

                4. **Node Data Format:**
                - `dataSource` Node:
                    - `data` contains the following properties:
                        {{
                            "prompt": "<data source description>",
                            "name": "{web_url}", !!!DOESN"T CHANGES 
                            "content": "<file name or content>",
                            "flow_id": "{flow_id}",
                            "file": "{web_url}"  // Empty object or file metadata
                        }}
                5. **Question Data Format:**
                - `question` Node:
                    - `data` contains the following properties:
                        {{
                            "question": "<the question asked for the response>",
                            "component_id": "<component reference ID - unique identifier of 12 or 24 digit unique uuid or nanoid>",
                            "component_type" : "web",
                        }}
                6. **RESPONSE NODE FORMAT**
                - `response` Node:
                    - `data` contains nested properties:
                        {{
                            "id": "<unique identifier of 12 or 24 digit unique uuid or nanoid>",
                            "type": "response" !!DOESN'T CHANGE,
                            "data": {{
                                "question": "<question text, if applicable>",
                                "summ": "<!!give me a detailed answer for the above question>",
                                "df": [],
                                "graph": "",
                                "flow_id": "{flow_id}",
                                "component_id": "<component reference ID - unique identifier of 12 or 24 digit unique uuid or nanoid>",
                                "component_type": "web"
                            }}
                        }}

                7. **Connections:**
                - Connections between nodes should be represented by edges, with the following format:
                    - `id` (unique identifier for the edge)
                    - `source` (ID of the source node)
                    - `target` (ID of the target node)
                    - `type` (optional, defaults to `default`)
                    - 'animated' !!WILL ALWAYS BE TRUE

                8. **Viewport Configuration:**
                - Include a `viewport` object that specifies:
                    - `x` (horizontal position of the viewport)
                    - `y` (vertical position of the viewport)
                    - `zoom` (zoom level for initial rendering)

                ### Additional Considerations:
                - Ensure that the node positions are distributed properly to avoid overlap.
                - Prioritize connecting `response` nodes where it adds logical structure to the flow.

                ### IMPORTANT:
                - **RETURN ONLY THE VALID JSON OBJECT AND NO ADDITIONAL COMMENTS**.
                - Do **not** include any explanations, text, or additional information.
                - Maintain the format with double curly braces `{{` and `}}` as shown in the format.
                """,   

                "attachments": [
                    {"file_id": messages_file.id, "tools": [{"type": "file_search"}]}
                ],
            }
            ]
        )

        run = openai.beta.threads.runs.create_and_poll(thread_id=thread.id, assistant_id=assistant.id)

        messages = list(openai.beta.threads.messages.list(thread_id=thread.id, run_id=run.id))
        
        message_content = messages[0].content[0].text
        annotations = message_content.annotations

        for index, annotation in enumerate(annotations):
            message_content.value = message_content.value.replace(annotation.text, f"[{index}]")

        response_json = message_content.value.replace("```json", "").replace("```", "").replace("\n", "").strip()
        
        response_json = json.loads(response_json)
        print(response_json)

        component_metadata = {
            "flow_id": ObjectId(flow_id),
            "web_url": web_url,
            "file_id": messages_file.id,
            "assistant_id": assistant.id,
            "vector_store_id": vector_store.id,
            "size": len(file_bytes),
            "type": "web",
            "processing_type": "gpt",
            "mindmap_json": response_json,
        }

        component_id = component_collection.insert_one(component_metadata).inserted_id

        return {
            "component_id": str(component_id),
            "type": "web",
            "mindmap_json": response_json,
            "flow_type": flow_type
        }

    except Exception as e:
        print(f"Error in /component-create-crawl endpoint: {e}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail="Internal Server Error")


@app.post("/pdf-component-qa", response_model=List[PDFNodeQueryResponse])
def PDF_QA(request: PDFNodeQueryRequest):
    try:
        record = component_collection.find_one(
            {
                "flow_id": ObjectId(request.flow_id),
                "_id": ObjectId(request.component_id),
                "type": "pdf",
            }
        )

        if not record or "processing_type" not in record:
            raise HTTPException(
                status_code=404,
                detail="processing_type not found for the given flow_id and component_id",
            )

        processing_type = record["processing_type"]

        if processing_type == "gpt":
            vector_store_id = record["vector_store_id"]
            file_id = record["file_id"]
            assistant_id = record["assistant_id"]
            response = one_shot_openai(
                request.query, vector_store_id, file_id, assistant_id
            )
            response_json = json.loads(response)
            print(response_json)

            node_data = {
                "_id": ObjectId(request.node_id),
                "flow_id": ObjectId(request.flow_id),
                "component_id": ObjectId(request.component_id),
                "question": request.query,
                "summ": response_json.get("summ", ""),
                "df": validate_dataframe(response_json.get("df", [])),
                "graph": response_json.get("graph", ""),
                "type": "pdf",
                "is_delete": "false",
                "timestamp": datetime.datetime.utcnow(),
            }

            node_id_response = node_collection.insert_one(node_data)
            
            print(node_id_response)
            
            question_entries = []

            question_entries.append(
                PDFNodeQueryResponse(
                    data={
                        "question": request.query,
                        "summ": response_json.get("summ", ""),
                        "df": validate_dataframe(response_json.get("df", [])),
                        "graph": response_json.get("graph", ""),
                        "flow_id": request.flow_id,
                        "component_id": request.component_id,
                        "component_type": "pdf",
                    },
                    id=request.node_id,
                    type="PDFNode",
                )
            )

            if request.request_type == "question":
                empty_question_entry = PDFNodeQueryResponse(
                    id=str(ObjectId()),
                    data={
                        "question": "",
                        "flow_id": request.flow_id,
                        "component_id": request.component_id,
                        "component_type": "pdf",
                    },
                    type="question",
                )

                question_entries.append(empty_question_entry)
                print(question_entries)
            return question_entries

        else:
            passages = get_relevant_passage(
                request.query, request.flow_id, request.component_id, 2
            )
            if not passages:
                raise HTTPException(
                    status_code=404, detail="No relevant passages found for the query."
                )

            relevant_passage = " ".join(passages)

            instructions = record["instructions"]

            template = """
                You are an AI assistant tasked with answering the user’s question based on the provided passages and the given persona. Return the results in **JSON format** with the structure below:  

                #### **Response Format:**  
                {{
                "summ": "Your summarized response here...",
                "df": an array of JSON objects,
                "graph": "json_string_representation_of_plotly_graph"
                }}

                ### **Instructions:**
                1. Answer the question using the passage.
                2. Extract relevant tabular data into a JSON object compatible with Ag-Grid. If no table exists, return empty JSON object.
                3. If a dataframe is available, generate a relevant **Plotly graph**. Return it as a **valid JSON string** that can be parsed in React.js.
                4. If no graph is possible, return an empty string `""`.
                5. ** The graph's background will be black, so adjust the theme accordingly**.

                NOTE -- "Make sure you need to return only json as response only & please don't add any comments"
                NOTE -- "Make sure you need only need the answer for which context of data is available if not available return empty json as per format"


                **Here is the question:** {query}  
                **Here is the persona:** {instructions}
                **Here is the passage: {escaped_passage}

                ### **Example Output:**  
                If the passage contains a table and a relevant graph, return:  

                ```json
                {{
                "summ": "Based on the passage, the key points discussed were...",
                "df": [
                    {{
                    "column1": "value1",
                    "column2": "value2",
                    "column3": "value3"
                    }},
                    {{
                    "column1": "value1",
                    "column2": "value2",
                    "column3": "value3"
                    }}
                ],
                "graph": "{{\"data\": [{{\"x\": [\"2024-02-01\", \"2024-02-02\"], \"y\": [100, 150], \"type\": \"line\"}}], \"layout\": {{\"title\": \"Sample Graph\"}}"
                }}
                }}
                """

            prompt = PromptTemplate.from_template(template)

            print(prompt)

            llm_chain = prompt | llm
            answer = llm_chain.invoke(
                {
                    "instructions": instructions,
                    "query": request.query,
                    "escaped_passage": relevant_passage,
                }
            )

            answer = answer.content.replace("```json", "").replace("```", "").strip()
            response = answer.replace("\n", "")
            response_json = json.loads(response)

            node_data = {
                "_id": ObjectId(request.node_id),
                "flow_id": ObjectId(request.flow_id),
                "component_id": ObjectId(request.component_id),
                "question": request.query,
                "summ": response_json.get("summ", ""),
                "df": validate_dataframe(response_json.get("df", [])),
                "graph": response_json.get("graph", ""),
                "type": "pdf",
                "is_delete": "false",
                "timestamp": datetime.datetime.utcnow(),
            }
            node_id_response = node_collection.insert_one(node_data)

            question_entries = []

            question_entries.append(
                PDFNodeQueryResponse(
                    data={
                        "question": request.query,
                        "summ": response_json.get("summ", ""),
                        "df": validate_dataframe(response_json.get("df", [])),
                        "graph": response_json.get("graph", ""),
                        "flow_id": request.flow_id,
                        "component_id": request.component_id,
                        "component_type": "pdf",
                    },
                    id=request.node_id,
                    type="PDFNode",
                )
            )

            if request.request_type == "question":
                empty_question_entry = PDFNodeQueryResponse(
                    id=str(ObjectId()),
                    data={
                        "question": "",
                        "flow_id": request.flow_id,
                        "component_id": request.component_id,
                        "component_type": "pdf",
                    },
                    type="question",
                )

                question_entries.append(empty_question_entry)

            return question_entries

    except Exception as e:
        print(traceback.print_exc())
        print(f"Error in /pdf-component-qa endpoint: {e.__traceback__}")
        raise HTTPException(status_code=500, detail="Internal Server Error")


@app.post("/txt-component-qa", response_model=List[TXTNodeQueryResponse])
def TXT_QA(request: TXTNodeQueryRequest):
    try:
        record = component_collection.find_one(
            {
                "flow_id": ObjectId(request.flow_id),
                "_id": ObjectId(request.component_id),
                "type": "txt",
            }
        )

        vector_store_id = record["vector_store_id"]
        file_id = record["file_id"]
        assistant_id = record["assistant_id"]
        response = one_shot_openai(
            request.query, vector_store_id, file_id, assistant_id
        )
        response_json = json.loads(response)
        print(response_json)

        node_data = {
            "_id": ObjectId(request.node_id),
            "flow_id": ObjectId(request.flow_id),
            "component_id": ObjectId(request.component_id),
            "question": request.query,
            "summ": response_json.get("summ", ""),
            "df": validate_dataframe(response_json.get("df", [])),
            "graph": response_json.get("graph", ""),
            "type": "txt",
            "is_delete": "false",
            "timestamp": datetime.datetime.utcnow(),
        }

        node_id_response = node_collection.insert_one(node_data)

        question_entries = []

        question_entries.append(
            TXTNodeQueryResponse(
                data={
                    "question": request.query,
                    "summ": response_json.get("summ", ""),
                    "df": validate_dataframe(response_json.get("df", [])),
                    "graph": response_json.get("graph", ""),
                    "flow_id": request.flow_id,
                    "component_id": request.component_id,
                    "component_type": "txt",
                },
                id=request.node_id,
                type="TXTNode",
            )
        )

        if request.request_type == "question":
            empty_question_entry = TXTNodeQueryResponse(
                id=str(ObjectId()),
                data={
                    "question": "",
                    "flow_id": request.flow_id,
                    "component_id": request.component_id,
                    "component_type": "txt",
                },
                type="question",
            )

            question_entries.append(empty_question_entry)
        return question_entries

    except Exception as e:
        print(f"Error in /txt-component-qa endpoint: {e.__traceback__}")
        raise HTTPException(status_code=500, detail="Internal Server Error")


@app.post("/md-component-qa", response_model=List[MDNodeQueryResponse])
def MD_QA(request: MDNodeQueryRequest):
    try:
        record = component_collection.find_one(
            {
                "flow_id": ObjectId(request.flow_id),
                "_id": ObjectId(request.component_id),
                "type": "md",
            }
        )

        vector_store_id = record["vector_store_id"]
        file_id = record["file_id"]
        assistant_id = record["assistant_id"]
        response = one_shot_openai(
            request.query, vector_store_id, file_id, assistant_id
        )
        response_json = json.loads(response)
        print(response_json)

        node_data = {
            "_id": ObjectId(request.node_id),
            "flow_id": ObjectId(request.flow_id),
            "component_id": ObjectId(request.component_id),
            "question": request.query,
            "summ": response_json.get("summ", ""),
            "df": validate_dataframe(response_json.get("df", [])),
            "graph": response_json.get("graph", ""),
            "type": "md",
            "is_delete": "false",
            "timestamp": datetime.datetime.utcnow(),
        }

        node_id_response = node_collection.insert_one(node_data)

        question_entries = []

        question_entries.append(
            MDNodeQueryResponse(
                data={
                    "question": request.query,
                    "summ": response_json.get("summ", ""),
                    "df": validate_dataframe(response_json.get("df", [])),
                    "graph": response_json.get("graph", ""),
                    "flow_id": request.flow_id,
                    "component_id": request.component_id,
                    "component_type": "md",
                },
                id=request.node_id,
                type="MDNode",
            )
        )

        if request.request_type == "question":
            empty_question_entry = MDNodeQueryResponse(
                id=str(ObjectId()),
                data={
                    "question": "",
                    "flow_id": request.flow_id,
                    "component_id": request.component_id,
                    "component_type": "md",
                },
                type="question",
            )

            question_entries.append(empty_question_entry)
            return question_entries

        return question_entries

    except Exception as e:
        print(f"Error in /md-component-qa endpoint: {e.__traceback__}")
        raise HTTPException(status_code=500, detail="Internal Server Error")


@app.post("/html-component-qa", response_model=List[HTMLNodeQueryResponse])
def HTML_QA(request: HTMLNodeQueryRequest):
    try:
        record = component_collection.find_one(
            {
                "flow_id": ObjectId(request.flow_id),
                "_id": ObjectId(request.component_id),
                "type": "html",
            }
        )

        vector_store_id = record["vector_store_id"]
        file_id = record["file_id"]
        assistant_id = record["assistant_id"]
        response = one_shot_openai(
            request.query, vector_store_id, file_id, assistant_id
        )
        response_json = json.loads(response)
        print(response_json)

        node_data = {
            "_id": ObjectId(request.node_id),
            "flow_id": ObjectId(request.flow_id),
            "component_id": ObjectId(request.component_id),
            "question": request.query,
            "summ": response_json.get("summ", ""),
            "df": validate_dataframe(response_json.get("df", [])),
            "graph": response_json.get("graph", ""),
            "type": "html",
            "is_delete": "false",
            "timestamp": datetime.datetime.utcnow(),
        }

        node_id_response = node_collection.insert_one(node_data)

        question_entries = []

        question_entries.append(
            HTMLNodeQueryResponse(
                data={
                    "question": request.query,
                    "summ": response_json.get("summ", ""),
                    "df": validate_dataframe(response_json.get("df", [])),
                    "graph": response_json.get("graph", ""),
                    "flow_id": request.flow_id,
                    "component_id": request.component_id,
                    "component_type": "html",
                },
                id=request.node_id,
                type="HTMLNode",
            )
        )

        if request.request_type == "question":
            empty_question_entry = HTMLNodeQueryResponse(
                id=str(ObjectId()),
                data={
                    "question": "",
                    "flow_id": request.flow_id,
                    "component_id": request.component_id,
                    "component_type": "html",
                },
                type="question",
            )

            question_entries.append(empty_question_entry)
            return question_entries

        return question_entries

    except Exception as e:
        print(f"Error in /html-component-qa endpoint: {e.__traceback__}")
        raise HTTPException(status_code=500, detail="Internal Server Error")


@app.post("/docx-component-qa", response_model=List[DOCXNodeQueryResponse])
def DOCX_QA(request: DOCXNodeQueryRequest):
    try:
        record = component_collection.find_one(
            {
                "flow_id": ObjectId(request.flow_id),
                "_id": ObjectId(request.component_id),
                "type": "docx",
            }
        )

        vector_store_id = record["vector_store_id"]
        file_id = record["file_id"]
        assistant_id = record["assistant_id"]
        response = one_shot_openai(
            request.query, vector_store_id, file_id, assistant_id
        )
        response_json = json.loads(response)
        print(response_json)

        node_data = {
            "_id": ObjectId(request.node_id),
            "flow_id": ObjectId(request.flow_id),
            "component_id": ObjectId(request.component_id),
            "question": request.query,
            "summ": response_json.get("summ", ""),
            "df": validate_dataframe(response_json.get("df", [])),
            "graph": response_json.get("graph", ""),
            "type": "docx",
            "is_delete": "false",
            "timestamp": datetime.datetime.utcnow(),
        }

        node_id_response = node_collection.insert_one(node_data)

        question_entries = []

        question_entries.append(
            DOCXNodeQueryResponse(
                data={
                    "question": request.query,
                    "summ": response_json.get("summ", ""),
                    "df": validate_dataframe(response_json.get("df", [])),
                    "graph": response_json.get("graph", ""),
                    "flow_id": request.flow_id,
                    "component_id": request.component_id,
                    "component_type": "docx",
                },
                id=request.node_id,
                type="DOCXNode",
            )
        )

        if request.request_type == "question":
            empty_question_entry = DOCXNodeQueryResponse(
                id=str(ObjectId()),
                data={
                    "question": "",
                    "flow_id": request.flow_id,
                    "component_id": request.component_id,
                    "component_type": "docx",
                },
                type="question",
            )

            question_entries.append(empty_question_entry)
            return question_entries

        return question_entries

    except Exception as e:
        print(f"Error in /docx-component-qa endpoint: {e.__traceback__}")
        raise HTTPException(status_code=500, detail="Internal Server Error")


@app.post("/pptx-component-qa", response_model=List[PPTXNodeQueryResponse])
def PPTX_QA(request: PPTXNodeQueryRequest):
    try:
        record = component_collection.find_one(
            {
                "flow_id": ObjectId(request.flow_id),
                "_id": ObjectId(request.component_id),
                "type": "pptx",
            }
        )

        vector_store_id = record["vector_store_id"]
        file_id = record["file_id"]
        assistant_id = record["assistant_id"]
        response = one_shot_openai(
            request.query, vector_store_id, file_id, assistant_id
        )
        response_json = json.loads(response)
        print(response_json)

        node_data = {
            "_id": ObjectId(request.node_id),
            "flow_id": ObjectId(request.flow_id),
            "component_id": ObjectId(request.component_id),
            "question": request.query,
            "summ": response_json.get("summ", ""),
            "df": validate_dataframe(response_json.get("df", [])),
            "graph": response_json.get("graph", ""),
            "type": "pptx",
            "is_delete": "false",
            "timestamp": datetime.datetime.utcnow(),
        }

        node_id_response = node_collection.insert_one(node_data)

        question_entries = []

        question_entries.append(
            PPTXNodeQueryResponse(
                data={
                    "question": request.query,
                    "summ": response_json.get("summ", ""),
                    "df": validate_dataframe(response_json.get("df", [])),
                    "graph": response_json.get("graph", ""),
                    "flow_id": request.flow_id,
                    "component_id": request.component_id,
                    "component_type": "pptx",
                },
                id=request.node_id,
                type="PPTXNode",
            )
        )

        if request.request_type == "question":
            empty_question_entry = PPTXNodeQueryResponse(
                id=str(ObjectId()),
                data={
                    "question": "",
                    "flow_id": request.flow_id,
                    "component_id": request.component_id,
                    "component_type": "pptx",
                },
                type="question",
            )

            question_entries.append(empty_question_entry)
            return question_entries

        return question_entries

    except Exception as e:
        print(f"Error in /pptx-component-qa endpoint: {e.__traceback__}")
        raise HTTPException(status_code=500, detail="Internal Server Error")


@app.post("/csv-component-qa", response_model=List[CSVNodeQueryResponse])
def CSV_QA(request: CSVNodeQueryRequest):
    try:
        record = component_collection.find_one(
            {
                "flow_id": ObjectId(request.flow_id),
                "_id": ObjectId(request.component_id),
                "type": "csv",
            }
        )

        if not record or "table_name" not in record:
            # Raise an HTTP 404 error if the table is not found
            raise HTTPException(
                status_code=404,
                detail="Table not found for the given flow_id and component_id",
            )

        table_name = record["table_name"]
        question_with_table = f"question - {request.query} for table name: {table_name}"
        sqlQuery = csvBot.generate_sql(question_with_table)
        if csvBot.is_sql_valid(sqlQuery):
            runSQLDF = csvBot.run_sql(sqlQuery)
            summSQL = csvBot.generate_summary(question_with_table, runSQLDF)
            code = csvBot.generate_plotly_code(
                question=question_with_table,
                sql=sqlQuery,
                df_metadata=f"Running df.dtypes gives:\n {runSQLDF.dtypes}",
            )
            fig = csvBot.get_plotly_figure(plotly_code=code, df=runSQLDF)
            plotyGraph = fig.to_json()
            df_dict = runSQLDF.to_dict(orient="records")
            df_dict = [{str(k): v for k, v in row.items()} for row in df_dict]
            result_document = {
                "_id": ObjectId(request.node_id),
                "question": request.query,
                "query": sqlQuery,
                "df": df_dict,
                "summ": summSQL,
                "graph": plotyGraph,
                "flow_id": ObjectId(request.flow_id),
                "component_id": ObjectId(request.component_id),
                "type": "csv",
                "is_delete": "false",
                "created_at": datetime.datetime.utcnow(),
            }
            node_id_response = node_collection.insert_one(result_document)
            question_entries = []
            response_data = {
                "id": str(request.node_id),
                "type": "CSVNode",
                "data": {
                    "question": request.query,
                    "df": runSQLDF.to_dict(orient="records"),
                    "summ": summSQL,
                    "graph": plotyGraph,
                    "flow_id": request.flow_id,
                    "component_id": request.component_id,
                    "component_type": "csv",
                },
            }
            question_entries.append(response_data)
            if request.request_type == "question":
                empty_node = {
                    "id": str(ObjectId()),
                    "type": "question",
                    "data": {
                        "question": "",
                        "flow_id": request.flow_id,
                        "component_id": request.component_id,
                        "component_type": "csv",
                    },
                }
                question_entries.append(empty_node)

            return question_entries

        else:
            print("No answer found from llm")
            df_dict = pd.DataFrame().to_dict(orient="records")
            df_dict = [{str(k): v for k, v in row.items()} for row in df_dict]

            result_document = {
                "_id": ObjectId(request.node_id),
                "question": request.query,
                "df": df_dict,
                "summ": "",
                "graph": "",
                "flow_id": ObjectId(request.flow_id),
                "component_id": ObjectId(request.component_id),
                "type": "csv",
                "is_delete": "false",
            }

            node_id_response = node_collection.insert_one(result_document)

        question_entries = []

        response_data = {
            "id": str(request.node_id),
            "type": "CSVNode",
            "data": {
                "question": request.query,
                "df": pd.DataFrame().to_dict(orient="records"),
                "summ": "",
                "graph": "",
                "flow_id": request.flow_id,
                "component_id": request.component_id,
                "component_type": "csv",
            },
        }

        question_entries.append(response_data)

        if request.request_type == "question":
            empty_node = {
                "id": str(ObjectId()),
                "type": "question",
                "data": {
                    "question": "",
                    "flow_id": request.flow_id,
                    "component_id": request.component_id,
                    "component_type": "csv",
                },
            }
            question_entries.append(empty_node)

            return question_entries
    except Exception as e:
        print(f"Error in /csv-component-qa: {e}")
        raise HTTPException(status_code=500, detail="Internal Server Error")


@app.post("/web-component-qa", response_model=List[WebNodeQueryResponse])
def WEB_QA(request: WebNodeQueryRequest):
    try:
        record = component_collection.find_one(
            {
                "flow_id": ObjectId(request.flow_id),
                "_id": ObjectId(request.component_id),
                "type": "web",
            }
        )
        vector_store_id = record["vector_store_id"]
        file_id = record["file_id"]
        assistant_id = record["assistant_id"]
        response = one_shot_openai(
            request.query, vector_store_id, file_id, assistant_id
        )
        response_json = json.loads(response)
        print(response_json)
        
        node_data = {
            "_id": ObjectId(request.node_id),
            "flow_id": ObjectId(request.flow_id),
            "component_id": ObjectId(request.component_id),
            "question": request.query,
            "summ": response_json.get("summ", ""),
            "df": validate_dataframe(response_json.get("df", [])),
            "graph": response_json.get("graph", {}),
            "type": "web",
            "is_delete": "false",
            "timestamp": datetime.datetime.utcnow(),
        }

        node_id_response = node_collection.insert_one(node_data)

        question_entries = []

        question_entries.append(
            WebNodeQueryResponse(
                data={
                    "question": request.query,
                    "summ": response_json.get("summ", ""),
                    "df": validate_dataframe(response_json.get("df", [])),
                    "graph": response_json.get("graph", {}),
                    "flow_id": request.flow_id,
                    "component_id": request.component_id,
                    "component_type": "web",
                },
                id=request.node_id,
                type="WebNode",
            )
        )

        if request.request_type == "question":
            empty_question_entry = WebNodeQueryResponse(
                id=str(ObjectId()),
                data={
                    "question": "",
                    "flow_id": request.flow_id,
                    "component_id": request.component_id,
                    "component_type": "web",
                },
                type="question",
            )

            question_entries.append(empty_question_entry)

        return question_entries

    except Exception as e:
        print(f"Error in /web-component-qa endpoint: {e}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail="Internal Server Error")


@app.post("/img-component-qa", response_model=List[ImgNodeQueryResponse])
def IMG_QA(request: ImgNodeQueryRequest):
    try:
        record = component_collection.find_one(
            {
                "flow_id": ObjectId(request.flow_id),
                "_id": ObjectId(request.component_id),
                "type": "image",
            }
        )

        instructions = record["instructions"]
        persona_name = record["persona_name"]
        base64_image = record["base64_image"]
        mime_type = record["mime_type"]
        image_bytes = base64.b64decode(base64_image)

        image_part = {"mime_type": mime_type, "data": image_bytes}

        template = f"""
    You are an AI assistant tasked with answering the user’s question based on the provided question and persona. Return the results in **JSON format** with the structure below:  

    #### **Response Format:**  
    {{
    "summ": "Your summarized response here...",
    "df": an array of JSON objects,
    "graph": "json_string_representation_of_plotly_graph"
    }}

    ### **Instructions:**
    1. Answer the question using the conversation history.
    2. Extract relevant tabular data into a JSON object compatible with Ag-Grid. If no table exists, return empty JSON object.
    3. If a dataframe is available, generate a relevant **Plotly graph**. Return it as a **valid JSON string** that can be parsed in React.js.
    4. If no graph is possible, return an empty string `""`.
    5. ** The graph's background will be black, so adjust the theme accordingly**.

    NOTE -- "Make sure you need to return only json as response only & please don't add any comments"
    NOTE -- "Make sure you need only need the answer for which context of data is available if not available return empty json as per format"

    **Here is the question:** {request.query}  
    **Here is the persona:** {persona_name}

    ### **Example Output:**  
    If the conversation history contains a table and a relevant graph, return:  

    ```json
    {{
    "summ": "Based on the conversation, the key points discussed were...",
    "df": [
        {{
        "column1": "value1",
        "column2": "value2",
        "column3": "value3"
        }},
        {{
        "column1": "value1",
        "column2": "value2",
        "column3": "value3"
        }}
    ],
    "graph": "{{\"data\": [{{\"x\": [\"2024-02-01\", \"2024-02-02\"], \"y\": [100, 150], \"type\": \"line\"}}], \"layout\": {{\"title\": \"Sample Graph\"}}"
    }}
    """

        if use_bedrock_only:
            # Use Claude for image analysis
            responseList = analyze_image_with_claude(base64_image, template, mime_type)
        else:
            # Use Gemini for image analysis
            response = model.generate_content(contents=[template, image_part])
            responseList = response.text

        responseList = (
            responseList.replace("```json", "")
            .replace("```", "")
            .replace("\n", "")
            .strip()
        )

        print(responseList)

        response_json = json.loads(responseList)
        print(response_json)

        node_data = {
            "_id": ObjectId(request.node_id),
            "flow_id": ObjectId(request.flow_id),
            "component_id": ObjectId(request.component_id),
            "question": request.query,
            "summ": response_json.get("summ", ""),
            "df": validate_dataframe(response_json.get("df", [])),
            "graph": response_json.get("graph", {}),
            "type": "image",
            "is_delete": "false",
            "timestamp": datetime.datetime.utcnow(),
        }

        node_id_response = node_collection.insert_one(node_data)

        question_entries = []

        question_entries.append(
            ImgNodeQueryResponse(
                data={
                    "question": request.query,
                    "summ": response_json.get("summ", ""),
                    "df": validate_dataframe(response_json.get("df", [])),
                    "graph": response_json.get("graph", {}),
                    "flow_id": request.flow_id,
                    "component_id": request.component_id,
                    "component_type": "image",
                },
                id=request.node_id,
                type="ImageNode",
            )
        )

        if request.request_type == "question":
            empty_question_entry = ImgNodeQueryResponse(
                id=str(ObjectId()),
                data={
                    "question": "",
                    "flow_id": request.flow_id,
                    "component_id": request.component_id,
                    "component_type": "image",
                },
                type="question",
            )

            question_entries.append(empty_question_entry)

        return question_entries

    except Exception as e:
        print(f"Error in /img-component-qa endpoint: {e}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail="Internal Server Error")


@app.post("/audio-component-qa", response_model=List[AudioNodeQueryResponse])
def AUDIO_QA(request: AudioNodeQueryRequest):
    try:
        record = component_collection.find_one(
            {
                "flow_id": ObjectId(request.flow_id),
                "_id": ObjectId(request.component_id),
                "type": "audio",
            }
        )
        
        instructions = record["instructions"]
        persona_name = record["persona_name"]

        base64_audio = record["base64_audio"]
        mime_type = record["mime_type"]
        audio_bytes = base64.b64decode(base64_audio)

        audio_part = {"mime_type": mime_type, "data": audio_bytes}

        template = f"""You are an AI assistant tasked with answering the user’s question based on the provided question and persona. Return the results in **JSON format** with the structure below:  

    #### **Response Format:**  
    {{
    "summ": "Your summarized response here...",
    "df": an array of JSON objects,
    "graph": "json_string_representation_of_plotly_graph"
    }}

    ### **Instructions:**
    1. Answer the question using the conversation history.
    2. Extract relevant tabular data into a JSON object compatible with Ag-Grid. If no table exists, return empty JSON object.
    3. If a dataframe is available, generate a relevant **Plotly graph**. Return it as a **valid JSON string** that can be parsed in React.js.
    4. If no graph is possible, return an empty string `""`.
    5. ** The graph's background will be black, so adjust the theme accordingly**.

    NOTE -- "Make sure you need to return only json as response only & please don't add any comments"
    NOTE -- "Make sure you need only need the answer for which context of data is available if not available return empty json as per format"

    **Here is the question:** {request.query}  
    **Here is the persona:** {persona_name}

    ### **Example Output:**  
    If the conversation history contains a table and a relevant graph, return:  

    ```json
    {{
    "summ": "Based on the conversation, the key points discussed were...",
    "df": [
        {{
        "column1": "value1",
        "column2": "value2",
        "column3": "value3"
        }},
        {{
        "column1": "value1",
        "column2": "value2",
        "column3": "value3"
        }}
    ],
    "graph": "{{\"data\": [{{\"x\": [\"2024-02-01\", \"2024-02-02\"], \"y\": [100, 150], \"type\": \"line\"}}], \"layout\": {{\"title\": \"Sample Graph\"}}"
    }}
    """

        if use_bedrock_only:
            # For audio QA with Bedrock, we need to upload to S3 first
            # This is a simplified approach - in production you might want to cache this
            unique_id = str(uuid4())
            s3_key = f"temp/audio/{unique_id}.audio"
            upload_to_s3(audio_bytes, bucket_name, s3_key)
            s3_uri = f"s3://{bucket_name}/{s3_key}"
            
            responseList = analyze_content_with_pegasus(s3_uri, template)
        else:
            # Use Gemini for audio analysis
            response = model.generate_content(contents=[template, audio_part])
            responseList = response.text

        responseList = (
            responseList.replace("```json", "")
            .replace("```", "")
            .replace("\n", "")
            .strip()
        )

        response_json = json.loads(responseList)
        print(response_json)

        node_data = {
            "_id": ObjectId(request.node_id),
            "flow_id": ObjectId(request.flow_id),
            "component_id": ObjectId(request.component_id),
            "question": request.query,
            "summ": response_json.get("summ", ""),
            "df": validate_dataframe(response_json.get("df", [])),
            "graph": response_json.get("graph", {}),
            "type": "audio",
            "is_delete": "false",
            "timestamp": datetime.datetime.utcnow(),
        }

        node_id_response = node_collection.insert_one(node_data)

        question_entries = []

        question_entries.append(
            AudioNodeQueryResponse(
                data={
                    "question": request.query,
                    "summ": response_json.get("summ", ""),
                    "df": validate_dataframe(response_json.get("df", [])),
                    "graph": response_json.get("graph", {}),
                    "flow_id": request.flow_id,
                    "component_id": request.component_id,
                    "component_type": "audio",
                },
                id=request.node_id,
                type="AudioNode",
            )
        )

        if request.request_type == "question":
            empty_question_entry = AudioNodeQueryResponse(
                id=str(ObjectId()),
                data={
                    "question": "",
                    "flow_id": request.flow_id,
                    "component_id": request.component_id,
                    "component_type": "audio",
                },
                type="question",
            )

            question_entries.append(empty_question_entry)

        return question_entries

    except Exception as e:
        print(f"Error in /audio-component-qa endpoint: {e}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail="Internal Server Error")


@app.post("/youtube-component-qa", response_model=List[YoutubeNodeQueryResponse])
def YOUTUBE_QA(request: YoutubeNodeQueryRequest):
    try:
        record = component_collection.find_one(
            {
                "flow_id": ObjectId(request.flow_id),
                "_id": ObjectId(request.component_id),
                "type": "youtube",
            }
        )

        instructions = record["instructions"]
        persona_name = record["persona_name"]
        
        youtube_url = record["youtube_url"]
        mime_type = "video/*"

        template = f"""
      You are an AI assistant tasked with answering the user’s question based on the provided question and persona. Return the results in **JSON format** with the structure below:  

    #### **Response Format:**  
    {{
    "summ": "Your summarized response here...",
    "df": an array of JSON objects,
    "graph": "json_string_representation_of_plotly_graph"
    }}

    ### **Instructions:**
    1. Answer the question using the conversation history.
    2. Extract relevant tabular data into a JSON object compatible with Ag-Grid. If no table exists, return empty JSON object.
    3. If a dataframe is available, generate a relevant **Plotly graph**. Return it as a **valid JSON string** that can be parsed in React.js.
    4. If no graph is possible, return an empty string `""`.
    5. ** The graph's background will be black, so adjust the theme accordingly**.

    NOTE -- "Make sure you need to return only json as response only & please don't add any comments"
    NOTE -- "Make sure you need only need the answer for which context of data is available if not available return empty json as per format"

    **Here is the question:** {request.query}  
    **Here is the persona:** {persona_name}

    ### **Example Output:**  
    If the conversation history contains a table and a relevant graph, return:  

    ```json
    {{
    "summ": "Based on the conversation, the key points discussed were...",
    "df": [
        {{
        "column1": "value1",
        "column2": "value2",
        "column3": "value3"
        }},
        {{
        "column1": "value1",
        "column2": "value2",
        "column3": "value3"
        }}
    ],
    "graph": "{{\"data\": [{{\"x\": [\"2024-02-01\", \"2024-02-02\"], \"y\": [100, 150], \"type\": \"line\"}}], \"layout\": {{\"title\": \"Sample Graph\"}}"
    }}
    """

        if use_bedrock_only:
            # Use Pegasus for YouTube analysis
            responseList = analyze_content_with_pegasus(youtube_url, template)
        else:
            # Use Gemini for YouTube analysis
            response = model_vertexai.generate_content(
                contents=[template, Part.from_uri(youtube_url, mime_type)]
            )
            responseList = response.text

        responseList = (
            responseList.replace("```json", "")
            .replace("```", "")
            .replace("\n", "")
            .strip()
        )

        response_json = json.loads(responseList)
        print(response_json)

        node_data = {
            "_id": ObjectId(request.node_id),
            "flow_id": ObjectId(request.flow_id),
            "component_id": ObjectId(request.component_id),
            "question": request.query,
            "summ": response_json.get("summ", ""),
            "df": validate_dataframe(response_json.get("df", [])),
            "graph": response_json.get("graph", {}),
            "type": "youtube",
            "is_delete": "false",
            "timestamp": datetime.datetime.utcnow(),
        }

        node_id_response = node_collection.insert_one(node_data)

        question_entries = []

        question_entries.append(
            YoutubeNodeQueryResponse(
                data={
                    "question": request.query,
                    "summ": response_json.get("summ", ""),
                    "df": validate_dataframe(response_json.get("df", [])),
                    "graph": response_json.get("graph", {}),
                    "flow_id": request.flow_id,
                    "component_id": request.component_id,
                    "component_type": "youtube",
                },
                id=request.node_id,
                type="YoutubeNode",
            )
        )

        if request.request_type == "question":
            empty_question_entry = YoutubeNodeQueryResponse(
                id=str(ObjectId()),
                data={
                    "question": "",
                    "flow_id": request.flow_id,
                    "component_id": request.component_id,
                    "component_type": "youtube",
                },
                type="question",
            )

            question_entries.append(empty_question_entry)

        return question_entries

    except Exception as e:
        print(f"Error in /youtube-component-qa endpoint: {e}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail="Internal Server Error")


@app.post("/video-component-qa", response_model=List[VideoNodeQueryResponse])
def VIDEO_QA(request: VideoNodeQueryRequest):
    try:
        record = component_collection.find_one(
            {
                "flow_id": ObjectId(request.flow_id),
                "_id": ObjectId(request.component_id),
                "type": "video",
            }
        )

        instructions = record["instructions"]
        persona_name = record["persona_name"]
        
        video_url = record["video_url"]
        mime_type = record["mime_type"]

        template = f"""
          You are an AI assistant tasked with answering the user’s question based on the provided question and persona. Return the results in **JSON format** with the structure below:  

    #### **Response Format:**  
    {{
    "summ": "Your summarized response here...",
    "df": an array of JSON objects,
    "graph": "json_string_representation_of_plotly_graph"
    }}

    ### **Instructions:**
    1. Answer the question using the conversation history.
    2. Extract relevant tabular data into a JSON object compatible with Ag-Grid. If no table exists, return empty JSON object.
    3. If a dataframe is available, generate a relevant **Plotly graph**. Return it as a **valid JSON string** that can be parsed in React.js.
    4. If no graph is possible, return an empty string `""`.
    5. ** The graph's background will be black, so adjust the theme accordingly**.

    NOTE -- "Make sure you need to return only json as response only & please don't add any comments"
    NOTE -- "Make sure you need only need the answer for which context of data is available if not available return empty json as per format"

    **Here is the question:** {request.query}  
    **Here is the persona:** {persona_name}

    ### **Example Output:**  
    If the conversation history contains a table and a relevant graph, return:  

    ```json
    {{
    "summ": "Based on the conversation, the key points discussed were...",
    "df": [
        {{
        "column1": "value1",
        "column2": "value2",
        "column3": "value3"
        }},
        {{
        "column1": "value1",
        "column2": "value2",
        "column3": "value3"
        }}
    ],
    "graph": "{{\"data\": [{{\"x\": [\"2024-02-01\", \"2024-02-02\"], \"y\": [100, 150], \"type\": \"line\"}}], \"layout\": {{\"title\": \"Sample Graph\"}}"
    }}
    """

        if use_bedrock_only:
            # Extract S3 URI from video_url for Pegasus
            # Assuming video_url is in format: https://bucket.s3.amazonaws.com/key
            if "s3.amazonaws.com" in video_url:
                # Convert HTTPS URL to S3 URI
                parts = video_url.split("/")
                bucket = parts[2].split(".")[0]
                key = "/".join(parts[3:])
                s3_uri = f"s3://{bucket}/{key}"
            else:
                s3_uri = video_url  # Assume it's already an S3 URI
            
            responseList = analyze_content_with_pegasus(s3_uri, template)
        else:
            # Use Gemini for video analysis
            response = model_vertexai.generate_content(
                contents=[template, Part.from_uri(video_url, mime_type)]
            )
            responseList = response.text

        responseList = (
            responseList.replace("```json", "")
            .replace("```", "")
            .replace("\n", "")
            .strip()
        )

        response_json = json.loads(responseList)
        print(response_json)

        node_data = {
            "_id": ObjectId(request.node_id),
            "flow_id": ObjectId(request.flow_id),
            "component_id": ObjectId(request.component_id),
            "question": request.query,
            "summ": response_json.get("summ", ""),
            "df": validate_dataframe(response_json.get("df", [])),
            "graph": response_json.get("graph", {}),
            "type": "video",
            "is_delete": "false",
            "timestamp": datetime.datetime.utcnow(),
        }

        node_id_response = node_collection.insert_one(node_data)

        question_entries = []

        question_entries.append(
            VideoNodeQueryResponse(
                data={
                    "question": request.query,
                    "summ": response_json.get("summ", ""),
                    "df": validate_dataframe(response_json.get("df", [])),
                    "graph": response_json.get("graph", {}),
                    "flow_id": request.flow_id,
                    "component_id": request.component_id,
                    "component_type": "video",
                },
                id=request.node_id,
                type="VideoNode",
            )
        )

        if request.request_type == "question":
            empty_question_entry = VideoNodeQueryResponse(
                id=str(ObjectId()),
                data={
                    "question": "",
                    "flow_id": request.flow_id,
                    "component_id": request.component_id,
                    "component_type": "video",
                },
                type="question",
            )

            question_entries.append(empty_question_entry)

        return question_entries

    except Exception as e:
        print(f"Error in /video-component-qa endpoint: {e}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail="Internal Server Error")


@app.delete("/soft-delete-node/{node_id}", response_model=dict)
def soft_delete_node(node_id: str):
    try:
        result = node_collection.update_one(
            {"_id": ObjectId(node_id)}, {"$set": {"is_delete": "true"}}
        )

        if result.matched_count == 0:
            raise HTTPException(status_code=404, detail="Node not found.")

        return {"message": "Node soft deleted successfully."}

    except Exception as e:
        print(f"Error in /soft-delete-node endpoint: {e}")
        raise HTTPException(status_code=500, detail="Internal Server Error")


@app.get("/get-all-flow-details/{flow_id}")
def get_flow_details(flow_id: str):
    try:
        if not ObjectId.is_valid(flow_id):
            raise HTTPException(status_code=400, detail="Invalid flow_id format.")

        flow = flow_collection.find_one({"_id": ObjectId(flow_id)})
        if not flow:
            raise HTTPException(status_code=404, detail="Flow not found.")

        components = list(component_collection.find({"flow_id": ObjectId(flow_id)}))
        if not components:
            raise HTTPException(
                status_code=404, detail="No components found for the given flow_id."
            )

        flow_details = {
            "flow": {
                "flow_id": str(flow["_id"]),
                "flow_name": flow.get("flow_name"),
                "description": flow.get("description"),
                "summary": flow.get("summary"),
            },
            "components": [],
        }

        for component in components:
            component_id = str(component["_id"])
            nodes = list(
                node_collection.find(
                    {
                        "component_id": ObjectId(component_id),
                        "flow_id": ObjectId(flow_id),
                    }
                )
            )

            component_details = {
                "component_id": component_id,
                "name": component.get("name"),
                "file_hash": component.get("file_hash"),
                "size": component.get("size"),
                "s3_path": component.get("s3_path"),
                "nodes": [],
            }

            for node in nodes:
                component_details["nodes"].append(
                    {
                        "node_id": str(node["_id"]),
                        "question": node.get("question"),
                        "answer": node.get("answer"),
                        "timestamp": node.get("timestamp"),
                    }
                )

            flow_details["components"].append(component_details)

        return flow_details

    except Exception as e:
        print(f"Error in /get-all-flow-details/{flow_id}: {e}")
        raise HTTPException(status_code=500, detail="Internal Server Error")


@app.post("/create_sql_component/", response_model=SQLComponentResponse)
def create_sql_component(request: SQLComponentRequest):
    try:
        flow = flow_collection.find_one({"_id": ObjectId(request.flow_id)})
        
        if flow["flow_type"] != 'manual':
            raise HTTPException(status_code=400, detail="Only Manual Mindmap is supported for SQL.")
        
        df_ddl = sqlBot.run_sql("SELECT type, sql FROM sqlite_master WHERE sql IS NOT NULL AND name LIKE '%" + request.table_name + "%'")
        print(df_ddl)

        for ddl in df_ddl['sql'].to_list():
            sqlBot.train(ddl=ddl)

        training_data = sqlBot.get_training_data()
        print(training_data)

        component_data = {
            "flow_id": ObjectId(request.flow_id),
            "type": "sql",
            "table_name": request.table_name,
        }

        component_id = component_collection.insert_one(component_data).inserted_id

        return SQLComponentResponse(
            component_id=str(component_id),
            type="sql",
            message="Component created successfully",
        )

    except Exception as e:
        print(f"Error in /create_sql_component/: {e}")
        raise HTTPException(
            status_code=500, detail=f"Error creating SQL component: {str(e)}"
        )


@app.post("/components-follow-up-questions", response_model=List[ComponentFollowUpQueryResponse])
def create_follow_up_questions(request: ComponentFollowUpQueryRequest):
    try:
        if request.component_type == "pdf":
            record = component_collection.find_one(
                {
                    "flow_id": ObjectId(request.flow_id),
                    "_id": ObjectId(request.component_id),
                    "type": "pdf",
                }
            )

            if record["processing_type"] == "gpt":
                assistant = openai.beta.assistants.update(
                    assistant_id=record["assistant_id"],
                    name=request.persona_name,
                    instructions=request.instructions,
                    model=request.model_name,
                    temperature=request.temperature,
                    top_p=request.top_p,
                    tool_resources={
                        "file_search": {"vector_store_ids": [record["vector_store_id"]]}
                    },
                )
            else:
                component_collection.update_one(
                    {"_id": ObjectId(request.component_id)},
                    {
                        "$set": {
                            "instructions": request.instructions,
                            "persona_name": request.persona_name,
                        }
                    },
                )

            summary_pdf = record["summary"]
            relevant_passage = " ".join(summary_pdf)
            template = """Given the following summary and persona, generate three follow-up questions that the persona might ask about the text data. Respond with a list of questions, one per line, in Python string format delimited by |||. If no relevant questions are found, return an empty string.

            Here is the summary :- {summary_pdf}
            Here is the persona :- {persona}"""

            prompt = PromptTemplate.from_template(template)

            llm_chain = prompt | llm
            answer = llm_chain.invoke(
                {"summary_pdf": relevant_passage, "persona": request.persona_name}
            )

            responseList = answer.content

            print(responseList)

            responseList = (
                responseList.replace("```python", "").replace("```", "").strip()
            )
            responseList = responseList.replace("\n", "")

            responseList = [item.strip() for item in responseList.split("|||")]

            question_entries = []

            if responseList:

                for q in responseList:
                    question_entries.append(
                        ComponentFollowUpQueryResponse(
                            id=str(ObjectId()),
                            flow_id=request.flow_id,
                            data={
                                "question": q,
                                "component_id": request.component_id,
                                "component_type": request.component_type,
                            },
                            type="followUp",
                            position={"x": 0, "y": 0},
                        )
                    )

            empty_question_entry = ComponentFollowUpQueryResponse(
                id=str(ObjectId()),
                flow_id=request.flow_id,
                position={"x": 0, "y": 0},
                data={
                    "question": "",
                    "component_id": request.component_id,
                    "component_type": request.component_type,
                },
                type="question",
            )

            question_entries.append(empty_question_entry)

            return question_entries

        elif request.component_type == "txt":
            record = component_collection.find_one(
                {
                    "flow_id": ObjectId(request.flow_id),
                    "_id": ObjectId(request.component_id),
                    "type": "txt",
                }
            )

            assistant = openai.beta.assistants.update(
                assistant_id=record["assistant_id"],
                name=request.persona_name,
                instructions=request.instructions,
                model=request.model_name,
                temperature=request.temperature,
                top_p=request.top_p,
                tool_resources={
                    "file_search": {"vector_store_ids": [record["vector_store_id"]]}
                },
            )


            summary_txt = record["summary"]

            relevant_passage = " ".join(summary_txt)

            template = """Given the following summary and persona, generate three follow-up questions that the persona might ask about the text data. Respond with a list of questions, one per line, in Python string format delimited by |||. If no relevant questions are found, return an empty string.

            Here is the summary :- {summary_txt}
            Here is the persona :- {persona}"""

            prompt = PromptTemplate.from_template(template)

            llm_chain = prompt | llm
            answer = llm_chain.invoke(
                {"summary_txt": relevant_passage, "persona": request.persona_name}
            )

            responseList = answer.content

            print(responseList)

            responseList = (
                responseList.replace("```python", "").replace("```", "").strip()
            )
            responseList = responseList.replace("\n", "")

            responseList = [item.strip() for item in responseList.split("|||")]

            question_entries = []

            if responseList:

                for q in responseList:
                    question_entries.append(
                        ComponentFollowUpQueryResponse(
                            id=str(ObjectId()),
                            flow_id=request.flow_id,
                            data={
                                "question": q,
                                "component_id": request.component_id,
                                "component_type": request.component_type,
                            },
                            type="followUp",
                            position={"x": 0, "y": 0},
                        )
                    )

            empty_question_entry = ComponentFollowUpQueryResponse(
                id=str(ObjectId()),
                flow_id=request.flow_id,
                position={"x": 0, "y": 0},
                data={
                    "question": "",
                    "component_id": request.component_id,
                    "component_type": request.component_type,
                },
                type="question",
            )

            question_entries.append(empty_question_entry)

            return question_entries

        elif request.component_type == "md":
            record = component_collection.find_one(
                {
                    "flow_id": ObjectId(request.flow_id),
                    "_id": ObjectId(request.component_id),
                    "type": "md",
                }
            )

            assistant = openai.beta.assistants.update(
                assistant_id=record["assistant_id"],
                name=request.persona_name,
                instructions=request.instructions,
                model=request.model_name,
                temperature=request.temperature,
                top_p=request.top_p,
                tool_resources={
                    "file_search": {"vector_store_ids": [record["vector_store_id"]]}
                },
            )

            summary_md = record["summary"]

            relevant_passage = " ".join(summary_md)

            template = """Given the following summary and persona, generate three follow-up questions that the persona might ask about the text data. Respond with a list of questions, one per line, in Python string format delimited by |||. If no relevant questions are found, return an empty string.

            Here is the summary :- {summary_md}
            Here is the persona :- {persona}"""

            prompt = PromptTemplate.from_template(template)

            llm_chain = prompt | llm
            answer = llm_chain.invoke(
                {"summary_md": relevant_passage, "persona": request.persona_name}
            )

            responseList = answer.content

            print(responseList)

            responseList = (
                responseList.replace("```python", "").replace("```", "").strip()
            )
            responseList = responseList.replace("\n", "")

            responseList = [item.strip() for item in responseList.split("|||")]

            question_entries = []

            if responseList:

                for q in responseList:
                    question_entries.append(
                        ComponentFollowUpQueryResponse(
                            id=str(ObjectId()),
                            flow_id=request.flow_id,
                            data={
                                "question": q,
                                "component_id": request.component_id,
                                "component_type": request.component_type,
                            },
                            type="followUp",
                            position={"x": 0, "y": 0},
                        )
                    )

            empty_question_entry = ComponentFollowUpQueryResponse(
                id=str(ObjectId()),
                flow_id=request.flow_id,
                position={"x": 0, "y": 0},
                data={
                    "question": "",
                    "component_id": request.component_id,
                    "component_type": request.component_type,
                },
                type="question",
            )

            question_entries.append(empty_question_entry)

            return question_entries

        elif request.component_type == "html":
            record = component_collection.find_one(
                {
                    "flow_id": ObjectId(request.flow_id),
                    "_id": ObjectId(request.component_id),
                    "type": "html",
                }
            )

            assistant = openai.beta.assistants.update(
                assistant_id=record["assistant_id"],
                name=request.persona_name,
                instructions=request.instructions,
                model=request.model_name,
                temperature=request.temperature,
                top_p=request.top_p,
                tool_resources={
                    "file_search": {"vector_store_ids": [record["vector_store_id"]]}
                },
            )

            summary_html = record["summary"]

            relevant_passage = " ".join(summary_html)

            template = """Given the following summary and persona, generate three follow-up questions that the persona might ask about the text data. Respond with a list of questions, one per line, in Python string format delimited by |||. If no relevant questions are found, return an empty string.

            Here is the summary :- {summary_html}
            Here is the persona :- {persona}"""

            prompt = PromptTemplate.from_template(template)

            llm_chain = prompt | llm
            answer = llm_chain.invoke(
                {"summary_html": relevant_passage, "persona": request.persona_name}
            )

            responseList = answer.content

            print(responseList)

            responseList = (
                responseList.replace("```python", "").replace("```", "").strip()
            )
            responseList = responseList.replace("\n", "")

            responseList = [item.strip() for item in responseList.split("|||")]

            question_entries = []

            if responseList:

                for q in responseList:
                    question_entries.append(
                        ComponentFollowUpQueryResponse(
                            id=str(ObjectId()),
                            flow_id=request.flow_id,
                            data={
                                "question": q,
                                "component_id": request.component_id,
                                "component_type": request.component_type,
                            },
                            type="followUp",
                            position={"x": 0, "y": 0},
                        )
                    )

            empty_question_entry = ComponentFollowUpQueryResponse(
                id=str(ObjectId()),
                flow_id=request.flow_id,
                position={"x": 0, "y": 0},
                data={
                    "question": "",
                    "component_id": request.component_id,
                    "component_type": request.component_type,
                },
                type="question",
            )

            question_entries.append(empty_question_entry)

            return question_entries

        elif request.component_type == "docx":
            record = component_collection.find_one(
                {
                    "flow_id": ObjectId(request.flow_id),
                    "_id": ObjectId(request.component_id),
                    "type": "docx",
                }
            )

            assistant = openai.beta.assistants.update(
                assistant_id=record["assistant_id"],
                name=request.persona_name,
                instructions=request.instructions,
                model=request.model_name,
                temperature=request.temperature,
                top_p=request.top_p,
                tool_resources={
                    "file_search": {"vector_store_ids": [record["vector_store_id"]]}
                },
            )

            summary_docx = record["summary"]

            relevant_passage = " ".join(summary_docx)

            template = """Given the following summary and persona, generate three follow-up questions that the persona might ask about the text data. Respond with a list of questions, one per line, in Python string format delimited by |||. If no relevant questions are found, return an empty string.

            Here is the summary :- {summary_docx}
            Here is the persona :- {persona}"""

            prompt = PromptTemplate.from_template(template)

            llm_chain = prompt | llm
            answer = llm_chain.invoke(
                {"summary_docx": relevant_passage, "persona": request.persona_name}
            )

            responseList = answer.content

            print(responseList)

            responseList = (
                responseList.replace("```python", "").replace("```", "").strip()
            )
            responseList = responseList.replace("\n", "")

            responseList = [item.strip() for item in responseList.split("|||")]

            question_entries = []

            if responseList:

                for q in responseList:
                    question_entries.append(
                        ComponentFollowUpQueryResponse(
                            id=str(ObjectId()),
                            flow_id=request.flow_id,
                            data={
                                "question": q,
                                "component_id": request.component_id,
                                "component_type": request.component_type,
                            },
                            type="followUp",
                            position={"x": 0, "y": 0},
                        )
                    )

            empty_question_entry = ComponentFollowUpQueryResponse(
                id=str(ObjectId()),
                flow_id=request.flow_id,
                position={"x": 0, "y": 0},
                data={
                    "question": "",
                    "component_id": request.component_id,
                    "component_type": request.component_type,
                },
                type="question",
            )

            question_entries.append(empty_question_entry)

            return question_entries

        elif request.component_type == "pptx":
            record = component_collection.find_one(
                {
                    "flow_id": ObjectId(request.flow_id),
                    "_id": ObjectId(request.component_id),
                    "type": "pptx",
                }
            )

            assistant = openai.beta.assistants.update(
                assistant_id=record["assistant_id"],
                name=request.persona_name,
                instructions=request.instructions,
                model=request.model_name,
                temperature=request.temperature,
                top_p=request.top_p,
                tool_resources={
                    "file_search": {"vector_store_ids": [record["vector_store_id"]]}
                },
            )


            summary_pptx = record["summary"]

            relevant_passage = " ".join(summary_pptx)

            template = """Given the following summary and persona, generate three follow-up questions that the persona might ask about the text data. Respond with a list of questions, one per line, in Python string format delimited by |||. If no relevant questions are found, return an empty string.

            Here is the summary :- {summary_pptx}
            Here is the persona :- {persona}"""

            prompt = PromptTemplate.from_template(template)

            llm_chain = prompt | llm
            answer = llm_chain.invoke(
                {"summary_pptx": relevant_passage, "persona": request.persona_name}
            )

            responseList = answer.content

            print(responseList)

            responseList = (
                responseList.replace("```python", "").replace("```", "").strip()
            )
            responseList = responseList.replace("\n", "")

            responseList = [item.strip() for item in responseList.split("|||")]

            question_entries = []

            if responseList:

                for q in responseList:
                    question_entries.append(
                        ComponentFollowUpQueryResponse(
                            id=str(ObjectId()),
                            flow_id=request.flow_id,
                            data={
                                "question": q,
                                "component_id": request.component_id,
                                "component_type": request.component_type,
                            },
                            type="followUp",
                            position={"x": 0, "y": 0},
                        )
                    )

            empty_question_entry = ComponentFollowUpQueryResponse(
                id=str(ObjectId()),
                flow_id=request.flow_id,
                position={"x": 0, "y": 0},
                data={
                    "question": "",
                    "component_id": request.component_id,
                    "component_type": request.component_type,
                },
                type="question",
            )

            question_entries.append(empty_question_entry)

            return question_entries

        elif request.component_type == "sql":
            record = component_collection.find_one(
                {
                    "flow_id": ObjectId(request.flow_id),
                    "_id": ObjectId(request.component_id),
                    "type": "sql",
                }
            )
            table_name = record["table_name"]
            responseDDL = sqlBot.get_related_ddl(table_name)
            responseDOC = sqlBot.get_related_documentation(table_name)
            responseSimilarSQL = sqlBot.get_similar_question_sql(table_name)
            responseList = sqlBot.get_followup_questions_custom(
                table_name, responseSimilarSQL, responseDDL, responseDOC
            )
            responseList = (
                responseList.replace("```python", "").replace("```", "").strip()
            )
            responseList = responseList.replace("\n", "")  # Remove newlines
            responseList = responseList.replace("\\", "")  # Remove backslashes

            responseList = [item.strip() for item in responseList.split("|||")]

            print(responseList)

            question_entries = []

            if responseList:

                for q in responseList:
                    question_entries.append(
                        ComponentFollowUpQueryResponse(
                            id=str(ObjectId()),
                            flow_id=request.flow_id,
                            data={
                                "question": q,
                                "component_id": request.component_id,
                                "component_type": request.component_type,
                            },
                            type="followUp",
                            position={"x": 0, "y": 0},
                        )
                    )

            empty_question_entry = ComponentFollowUpQueryResponse(
                id=str(ObjectId()),
                flow_id=request.flow_id,
                position={"x": 0, "y": 0},
                data={
                    "question": "",
                    "component_id": request.component_id,
                    "component_type": request.component_type,
                },
                type="question",
            )

            question_entries.append(empty_question_entry)

            return question_entries

        elif request.component_type == "csv":
            record = component_collection.find_one(
                {
                    "flow_id": ObjectId(request.flow_id),
                    "_id": ObjectId(request.component_id),
                    "type": "csv",
                }
            )
            table_name = record["table_name"]
            print("Thissssss is table name", table_name)
            responseDDL = csvBot.get_related_ddl(table_name)
            print("THIS IS RESSSPONSE -----", responseDDL)
            responseDOC = csvBot.get_related_documentation(table_name)
            responseSimilarSQL = csvBot.get_similar_question_sql(table_name)
            responseList = csvBot.get_followup_questions_custom(
                table_name, responseSimilarSQL, responseDDL, responseDOC
            )
            responseList = (
                responseList.replace("```python", "").replace("```", "").strip()
            )
            responseList = responseList.replace("\n", "")
            responseList = responseList.replace("\\", "")  # Remove backslashes

            responseList = [item.strip() for item in responseList.split("|||")]

            print(responseList)

            question_entries = []

            if responseList:

                for q in responseList:
                    question_entries.append(
                        ComponentFollowUpQueryResponse(
                            id=str(ObjectId()),
                            flow_id=request.flow_id,
                            data={
                                "question": q,
                                "component_id": request.component_id,
                                "component_type": request.component_type,
                            },
                            type="followUp",
                            position={"x": 0, "y": 0},
                        )
                    )

            empty_question_entry = ComponentFollowUpQueryResponse(
                id=str(ObjectId()),
                flow_id=request.flow_id,
                position={"x": 0, "y": 0},
                data={
                    "question": "",
                    "component_id": request.component_id,
                    "component_type": request.component_type,
                },
                type="question",
            )

            question_entries.append(empty_question_entry)

            return question_entries

        elif request.component_type == "web":
            record = component_collection.find_one(
                {
                    "flow_id": ObjectId(request.flow_id),
                    "_id": ObjectId(request.component_id),
                    "type": "web",
                }
            )
            
            assistant = openai.beta.assistants.update(
                assistant_id=record["assistant_id"],
                name=request.persona_name,
                instructions=request.instructions,
                model=request.model_name,
                temperature=request.temperature,
                top_p=request.top_p,
                tool_resources={
                    "file_search": {"vector_store_ids": [record["vector_store_id"]]}
                },
            )

            summary_web = record["summary"]

            relevant_passage = " ".join(summary_web)

            template = """
            Given the following summary and persona, generate three follow-up questions that the persona might ask about the text data. Respond with a list of questions, one per line, in Python string format delimited by |||. If no relevant questions are found, return an empty string.

            Here is the summary :- {summary_web}
            Here is the persona :- {persona}
            """

            prompt = PromptTemplate.from_template(template)

            llm_chain = prompt | llm
            answer = llm_chain.invoke(
                {"summary_web": relevant_passage, "persona": request.persona_name}
            )

            responseList = answer.content

            print(responseList)

            responseList = (
                responseList.replace("```python", "").replace("```", "").strip()
            )
            responseList = responseList.replace("\n", "")  # Remove newlines

            responseList = [item.strip() for item in responseList.split("|||")]

            question_entries = []

            if responseList:

                for q in responseList:
                    question_entries.append(
                        ComponentFollowUpQueryResponse(
                            id=str(ObjectId()),
                            flow_id=request.flow_id,
                            data={
                                "question": q,
                                "component_id": request.component_id,
                                "component_type": request.component_type,
                            },
                            type="followUp",
                            position={"x": 0, "y": 0},
                        )
                    )

            empty_question_entry = ComponentFollowUpQueryResponse(
                id=str(ObjectId()),
                flow_id=request.flow_id,
                position={"x": 0, "y": 0},
                data={
                    "question": "",
                    "component_id": request.component_id,
                    "component_type": request.component_type,
                },
                type="question",
            )

            question_entries.append(empty_question_entry)

            return question_entries

        elif request.component_type == "image":
            record = component_collection.find_one(
                {
                    "flow_id": ObjectId(request.flow_id),
                    "_id": ObjectId(request.component_id),
                    "type": "image",
                }
            )

            if record:
                component_collection.update_one(
                    {"_id": ObjectId(request.component_id)},
                    {
                        "$set": {
                            "instructions": request.instructions,
                            "persona_name": request.persona_name,
                        }
                    },
                )

                base64_image = record["base64_image"]
                mime_type = record["mime_type"]
                image_bytes = base64.b64decode(base64_image)

                image_part = {"mime_type": mime_type, "data": image_bytes}
                #    and persona - "+request.persona_name+" ,

                response = model.generate_content(
                    contents=[
                        "Given the following image generate three follow-up questions that the persona - "+ request.persona_name +" might ask about the image. Respond with a list of questions, one per line, in Python string format delimited by |||. If no relevant questions are found, return an empty string and return only questions only",
                        image_part,
                    ]
                )

                responseList = response.text
                responseList = (
                    responseList.replace("```python", "").replace("```", "").strip()
                )
                responseList = [
                    q.strip()
                    for q in responseList.strip('"').replace("\n", "").split("|||")
                ]

                print(responseList)

                question_entries = []

                if responseList:

                    for q in responseList:
                        question_entries.append(
                            ComponentFollowUpQueryResponse(
                                id=str(ObjectId()),
                                flow_id=request.flow_id,
                                data={
                                    "question": q,
                                    "component_id": request.component_id,
                                    "component_type": request.component_type,
                                },
                                type="followUp",
                                position={"x": 0, "y": 0},
                            )
                        )

                empty_question_entry = ComponentFollowUpQueryResponse(
                    id=str(ObjectId()),
                    flow_id=request.flow_id,
                    position={"x": 0, "y": 0},
                    data={
                        "question": "",
                        "component_id": request.component_id,
                        "component_type": request.component_type,
                    },
                    type="question",
                )
                
                question_entries.append(empty_question_entry)

                return question_entries
        elif request.component_type == "audio":
            record = component_collection.find_one(
                {
                    "flow_id": ObjectId(request.flow_id),
                    "_id": ObjectId(request.component_id),
                    "type": "audio",
                }
            )

            if record:

                component_collection.update_one(
                    {"_id": ObjectId(request.component_id)},
                    {
                        "$set": {
                            "instructions": request.instructions,
                            "persona_name": request.persona_name,
                        }
                    },
                )

                base64_audio = record["base64_audio"]
                mime_type = record["mime_type"]
                audio_bytes = base64.b64decode(base64_audio)

                audio_part = {"mime_type": mime_type, "data": audio_bytes}

                response = model.generate_content(
                    contents=[
                        "Given the following audio generate three follow-up questions that the persona - "+ request.persona_name +" might ask about the image. Respond with a list of questions, one per line, in Python string format delimited by |||. If no relevant questions are found, return an empty string and return only questions only",
                        audio_part,
                    ]
                )

                responseList = response.text
                responseList = (
                    responseList.replace("```python", "").replace("```", "").strip()
                )
                responseList = [
                    q.strip()
                    for q in responseList.strip('"').replace("\n", "").split("|||")
                ]

                print(responseList)

                question_entries = []

                if responseList:

                    for q in responseList:
                        question_entries.append(
                            ComponentFollowUpQueryResponse(
                                id=str(ObjectId()),
                                flow_id=request.flow_id,
                                data={
                                    "question": q,
                                    "component_id": request.component_id,
                                    "component_type": request.component_type,
                                },
                                type="followUp",
                                position={"x": 0, "y": 0},
                            )
                        )

                empty_question_entry = ComponentFollowUpQueryResponse(
                    id=str(ObjectId()),
                    flow_id=request.flow_id,
                    position={"x": 0, "y": 0},
                    data={
                        "question": "",
                        "component_id": request.component_id,
                        "component_type": request.component_type,
                    },
                    type="question",
                )
                
                question_entries.append(empty_question_entry)

                return question_entries

        elif request.component_type == "youtube":
            record = component_collection.find_one(
                {
                    "flow_id": ObjectId(request.flow_id),
                    "_id": ObjectId(request.component_id),
                    "type": "youtube",
                }
            )

            if record:

                component_collection.update_one(
                    {"_id": ObjectId(request.component_id)},
                    {
                        "$set": {
                            "instructions": request.instructions,
                            "persona_name": request.persona_name,
                        }
                    },
                )

                youtube_url = record["youtube_url"]
                mime_type = "video/*"

                response = model_vertexai.generate_content(
                    contents=[
                        "Given the following youtube video generate three follow-up questions that the persona - "+ request.persona_name +" might ask about the image. Respond with a list of questions, one per line, in Python string format delimited by |||. If no relevant questions are found, return an empty string and return only questions only",
                        Part.from_uri(youtube_url, mime_type),
                    ]
                )
                
                responseList = response.text
                responseList = (
                    responseList.replace("```python", "").replace("```", "").strip()
                )
                responseList = [
                    q.strip()
                    for q in responseList.strip('"').replace("\n", "").split("|||")
                ]

                print(responseList)

                question_entries = []

                if responseList:

                    for q in responseList:
                        question_entries.append(
                            ComponentFollowUpQueryResponse(
                                id=str(ObjectId()),
                                flow_id=request.flow_id,
                                data={
                                    "question": q,
                                    "component_id": request.component_id,
                                    "component_type": request.component_type,
                                },
                                type="followUp",
                                position={"x": 0, "y": 0},
                            )
                        )

                empty_question_entry = ComponentFollowUpQueryResponse(
                    id=str(ObjectId()),
                    flow_id=request.flow_id,
                    position={"x": 0, "y": 0},
                    data={
                        "question": "",
                        "component_id": request.component_id,
                        "component_type": request.component_type,
                    },
                    type="question",
                )

                question_entries.append(empty_question_entry)

                return question_entries

        elif request.component_type == "video":
            record = component_collection.find_one(
                {
                    "flow_id": ObjectId(request.flow_id),
                    "_id": ObjectId(request.component_id),
                    "type": "video",
                }
            )

            if record:

                component_collection.update_one(
                    {"_id": ObjectId(request.component_id)},
                    {
                        "$set": {
                            "instructions": request.instructions,
                            "persona_name": request.persona_name,
                        }
                    },
                )

                video_url = record["video_url"]
                mime_type = record["mime_type"]


                response = model_vertexai.generate_content(
                    contents=[
                        "Given the following video generate three follow-up questions that the persona - "+ request.persona_name +" might ask about the image. Respond with a list of questions, one per line, in Python string format delimited by |||. If no relevant questions are found, return an empty string and return only questions only",
                        Part.from_uri(video_url, mime_type),
                    ]
                )

                responseList = response.text
                responseList = (
                    responseList.replace("```python", "").replace("```", "").strip()
                )
                responseList = [
                    q.strip()
                    for q in responseList.strip('"').replace("\n", "").split("|||")
                ]

                question_entries = []

                if responseList:

                    for q in responseList:
                        question_entries.append(
                            ComponentFollowUpQueryResponse(
                                id=str(ObjectId()),
                                flow_id=request.flow_id,
                                data={
                                    "question": q,
                                    "component_id": request.component_id,
                                    "component_type": request.component_type,
                                },
                                type="followUp",
                                position={"x": 0, "y": 0},
                            )
                        )

                empty_question_entry = ComponentFollowUpQueryResponse(
                    id=str(ObjectId()),
                    flow_id=request.flow_id,
                    position={"x": 0, "y": 0},
                    data={
                        "question": "",
                        "component_id": request.component_id,
                        "component_type": request.component_type,
                    },
                    type="question",
                )

                question_entries.append(empty_question_entry)

                return question_entries
        else:
            pass

    except Exception as e:
        print(f"Error in /components-follow-up-questions: {e}")
        raise HTTPException(status_code=500, detail="Internal Server Error")


@app.post("/sql-component-qa", response_model=List[SQLNodeQueryResponse])
def SQL_QA(request: SQLNodeQueryRequest):
    try:
        record = component_collection.find_one(
            {
                "flow_id": ObjectId(request.flow_id),
                "_id": ObjectId(request.component_id),
                "type": "sql",
            }
        )

        if not record or "table_name" not in record:
            raise HTTPException(
                status_code=404,
                detail="Table not found for the given flow_id and component_id",
            )

        table_name = record["table_name"]

        question_with_table = (
            f"question - {request.question} for table name: {table_name}"
        )

        sqlQuery = sqlBot.generate_sql(question_with_table)
        if sqlBot.is_sql_valid(sqlQuery):
            runSQLDF = sqlBot.run_sql(sqlQuery)
            summSQL = sqlBot.generate_summary(question_with_table, runSQLDF)
            code = sqlBot.generate_plotly_code(
                question=question_with_table,
                sql=sqlQuery,
                df_metadata=f"Running df.dtypes gives:\n {runSQLDF.dtypes}",
            )
            fig = sqlBot.get_plotly_figure(plotly_code=code, df=runSQLDF)
            plotyGraph = fig.to_json()

            df_dict = runSQLDF.to_dict(orient="records")
            df_dict = [{str(k): v for k, v in row.items()} for row in df_dict]
            result_document = {
                "_id": ObjectId(request.node_id),
                "question": request.question,
                "query": sqlQuery,
                "df": df_dict,
                "summ": summSQL,
                "graph": plotyGraph,
                "flow_id": ObjectId(request.flow_id),
                "component_id": ObjectId(request.component_id),
                "type": "sql",
                "is_delete": "false",
                "created_at": datetime.datetime.utcnow(),
            }

            node_id_response = node_collection.insert_one(result_document)

            question_entries = []

            response_data = {
                "id": str(request.node_id),
                "type": "SQLNode",
                "data": {
                    "question": request.question,
                    "query": sqlQuery,
                    "df": runSQLDF.to_dict(orient="records"),
                    "summ": summSQL,
                    "graph": plotyGraph,
                    "flow_id": request.flow_id,
                    "component_id": request.component_id,
                    "component_type": "sql",
                },
            }

            question_entries.append(response_data)

            if request.request_type == "question":
                empty_node = {
                    "id": str(ObjectId()),
                    "type": "question",
                    "data": {
                        "question": "",
                        "flow_id": request.flow_id,
                        "component_id": request.component_id,
                        "component_type": "sql",
                    },
                }
                question_entries.append(empty_node)

            return question_entries

        else:

            print("No answer found from llm")

            df_dict = pd.DataFrame().to_dict(orient="records")
            df_dict = [{str(k): v for k, v in row.items()} for row in df_dict]

            result_document = {
                "_id": ObjectId(request.node_id),
                "question": request.question,
                "query": "I don't know",
                "df": df_dict,
                "summ": "",
                "graph": "",
                "flow_id": ObjectId(request.flow_id),
                "component_id": ObjectId(request.component_id),
                "type": "sql",
                "is_delete": "false",
                "created_at": datetime.datetime.utcnow(),
            }

            node_id_response = node_collection.insert_one(result_document)

            question_entries = []

            response_data = {
                "id": str(request.node_id),
                "type": "SQLNode",
                "data": {
                    "question": request.question,
                    "query": "I don't know",
                    "df": pd.DataFrame().to_dict(orient="records"),
                    "summ": "",
                    "graph": "",
                    "flow_id": request.flow_id,
                    "component_id": request.component_id,
                    "component_type": "sql",
                },
            }

            question_entries.append(response_data)

            if request.request_type == "question":
                empty_node = {
                    "id": str(ObjectId()),
                    "type": "question",
                    "data": {
                        "question": "",
                        "flow_id": request.flow_id,
                        "component_id": request.component_id,
                        "component_type": "sql",
                    },
                }

                question_entries.append(empty_node)

            return question_entries

    except Exception as e:
        traceback.print_exc()
        print(f"Error in /sql-component-qa: {str(e.__traceback__)}")

        raise HTTPException(status_code=500, detail="Internal Server Error")


@app.post("/multiple-qa-summarize", response_model=MultipleQuestionAnswerQueryResponse)
def multiple_qa_summarize(request: MultipleQuestionAnswerQueryRequest):
    try:
        conversation = []

        for parent_id in request.parent_node_ids:
            question, answer = fetch_question_answer_from_node_collection(
                parent_id, request.flow_id
            )
            if question and answer:
                conversation.append({"role": "user", "content": question})
                conversation.append(
                    {"role": "assistant", "content": f"Answer: \n{answer}"}
                )

        print(conversation)

        template = """
        You are an AI assistant tasked with answering the user’s question based on the provided conversation history. Return the results in **JSON format** with the structure below:  

        #### **Response Format:**  
        {{
        "summ": "Your summarized response here...",
        "df": an array of JSON objects,
        "graph": "json_string_representation_of_plotly_graph"
        }}

        ### **Instructions:**
        1. Answer the question using the conversation history.
        2. Extract relevant tabular data into a JSON object compatible with Ag-Grid. If no table exists, return empty JSON object.
        3. If a dataframe is available, generate a relevant **Plotly graph**. Return it as a **valid JSON string** that can be parsed in React.js.
        4. If no graph is possible, return an empty string `""`.
        5. ** The graph's background will be black, so adjust the theme accordingly**.

        **Here is the question:** {query}  
        **Here is the conversation history:** {history}

        NOTE -- "Make sure you need to return only json as response only & please don't add any comments"
        NOTE -- "Make sure you need only need the answer for which context of data is available if not available return empty json as per format"
        

        ### **Example Output:**  
        If the conversation history contains a table and a relevant graph, return:  

        ```json
        {{
        "summ": "Based on the conversation, the key points discussed were...",
        "df": [
            {{
            "column1": "value1",
            "column2": "value2",
            "column3": "value3"
            }},
            {{
            "column1": "value1",
            "column2": "value2",
            "column3": "value3"
            }}
        ],
        "graph": "{{\"data\": [{{\"x\": [\"2024-02-01\", \"2024-02-02\"], \"y\": [100, 150], \"type\": \"line\"}}], \"layout\": {{\"title\": \"Sample Graph\"}}"
        }}
        """

        json_structure = """{
            "summ": "Your summary text here...",
                "df": [
                    {
                    "column1": "value1",
                    "column2": "value2",
                    "column3": "value3"
                    },
                    {
                    "column1": "value1",
                    "column2": "value2",
                    "column3": "value3"
                    }
                ],
                "graph": {
                    "data": [ ... ],
                    "layout": { ... }
                }
            }"""

        prompt = PromptTemplate.from_template(template)

        print(prompt)

        llm_chain = prompt | llm
        answer = llm_chain.invoke({"query": request.question, "history": conversation})

        response = (
            answer.content.replace("```json", "")
            .replace("```", "")
            .strip()
            .replace("\n", "")
        )
        response_json = json.loads(response)
        print(response_json)

        node_data = {
            "_id": ObjectId(request.node_id),
            "flow_id": ObjectId(request.flow_id),
            "parent_node_ids": request.parent_node_ids,
            "question": request.question,
            "summ": response_json.get("summ", ""),
            "df": validate_dataframe(response_json.get("df", [])),
            "graph": response_json.get("graph", ""),
            "type": "MultipleQA",
            "is_delete": "false",
            "created_at": datetime.datetime.utcnow(),
        }

        node_id_response = node_collection.insert_one(node_data)

        response = MultipleQuestionAnswerQueryResponse(
            data={
                "question": request.question,
                "summ": response_json.get("summ", ""),
                "df": validate_dataframe(response_json.get("df", [])),
                "graph": response_json.get("graph", ""),
                "flow_id": request.flow_id,
                "parent_node_ids": request.parent_node_ids,
                "component_type": "MultipleQA",
            },
            id=request.node_id,
            type="MultipleQA",
            parent_node_ids=request.parent_node_ids,
        )

        return response

    except Exception as e:
        print(f"Error in /multiple-qa-summarize: {e}")
        raise HTTPException(
            status_code=500, detail=f"Error processing the request: {str(e)}"
        )


@app.post("/flow-summarizer", response_model=FlowSummarizeResponse)
def flow_summarizer(request: FlowSummarizeRequest):
    try:
        nodes = node_collection.find({"flow_id": ObjectId(request.flow_id)})
        if not nodes:
            raise HTTPException(
                status_code=404, detail="No nodes found for the given flow_id."
            )

        conversation = []
        print("This are nodes", request.flow_id)
        for node in nodes:
            node_id = node["_id"]
            print(node_id)
            question, answer = fetch_question_answer_from_node_collection(
                node_id, request.flow_id
            )
            if question and answer:
                conversation.append({"role": "user", "content": question})
                conversation.append(
                    {"role": "assistant", "content": f"Answer: \n{answer}"}
                )

        print(conversation)

        template = """You are an AI assistant tasked with generating a JSX element based on multiple conversations between the user and assistant.

        ### **Rules:**
        1. **Return only valid JSX. No explanations, comments, or extra text.**
        2. **Use only `plotly.js` and `ag-grid-community`. No other libraries are allowed.**
        3. **The output must look like a structured financial report, not just 2-3 components.**
        4. **The layout should have clear sections like:**
            - Executive Summary (must be 100-150 words) 
            - Key Financial Metrics  (must highlight all the crucial points)
            - Performance Tables  (most important and data crucial)
            - Multiple Charts for Trends  
            - Additional Insights  
        5. **Ensure proper spacing, professional styling, and structured formatting.**
        6. **Use `ag-grid-community` for multiple tables.**
        7. **Use `plotly.js` for multiple relevant graphs.**
        8. **All sections should be visually distinct but cohesive.**
        9. **Make the design responsive, with professional styling (flexbox, grid, typography).**
        10. **Ensure the final output looks like an actual financial report from an investment firm.**
        11. For AG-Grid React always include the ref, rowClass, rowHeight, rowStyle, headerHeight and domLayout given in the **REFERENCE**.
        12. **FONT COLOR SHOULD BE BLACK FOR p tags and h1 tags**.
        13. **!important NO STYLING SHOULD BE APPLIED TO PLOT**.
        14. **BACKGROUND COLOR  OF MAIN DIV WILL BE WHITE**.
        15. **TRY TO LAYOUT AG-GRID BY GIVING HEIGHT AND WIDTH AS INLINE CSS BASED ON THE ROWDATA AND COLDEFS**.
        16. **!important GET AG-GRID TABLES IN WHITE **.
        17. **INCLUDE HR TAG after EACH SECTION**.
        18. **ALWAYS HAVE THE HEADING OF FINANCIAL REPORT AT BEGINNING AT THE CENTER**.
        19. **GIVE SOME GAPS AFTER EACH SECTION AND BETWEEN AG-GRID TABLES**.


        ### **Conversation History**
        Here is the conversation history :- {conversation}

        ### **Reference Output (Only JSX, No Comments or Extra Text, No Need to follow as given below):**  

        <div>
            <p style={{ fontSize: "16px", fontWeight: "bold", marginBottom: "10px" }}>
                Summary: {{Your Answer}}
            </p>

            <div className="ag-theme-alpine"> 
                <AgGridReact
                    rowData=[
                        {{ column1: "value1", column2: "value2", column3: "value3" }},
                        {{ column1: "value1", column2: "value2", column3: "value3" }}
                    ]
                    columnDefs=[
                        {{ headerName: "Column 1", field: "column1" }},
                        {{ headerName: "Column 2", field: "column2" }},
                        {{ headerName: "Column 3", field: "column3" }}
                    ]
                    rowClass={{"ag-row"}}
					rowHeight={{56}}
					rowStyle={{ alignItems: "center !important" }}
					headerHeight={{56}}
                    domLayout="autoHeight"
                />
            </div>

            <div style={{ width: "100%", height: "400px" }}>
                <Plot
                    data=[
                        {{ x: [1, 2, 3], y: [10, 20, 30], type: "scatter", mode: "lines+markers", marker: {{ color: "red" }}}}
                    ]
                    layout={{ title: "Graph Title", width: 600, height: 400 }}
                />
            </div>
        </div>        
        """
        prompt = PromptTemplate.from_template(template)

        llm_chain = prompt | llm
        answer = llm_chain.invoke({"conversation": conversation})

        print(answer)

        answer = answer.content.replace("```jsx", "").replace("```", "").strip()
        response = answer.replace("\n", "")

        update_result = flow_collection.update_one(
            {"_id": ObjectId(request.flow_id)}, {"$set": {"summary": response}}
        )

        response = FlowSummarizeResponse(
            flow_id=request.flow_id,
            response=response,
        )
        print(response)

        return response

    except Exception as e:
        print(f"Error in /flow-summarizer: {e}")
        traceback.print_exc()
        raise HTTPException(
            status_code=500, detail=f"Error processing the request: {str(e)}"
        )

@app.get("/sqlite-tables", response_model=List[str])
def read_sqlite_tables():
    cursor = connection.cursor()
    cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
    tables = [row[0] for row in cursor.fetchall()]
    connection.close()
    return tables
